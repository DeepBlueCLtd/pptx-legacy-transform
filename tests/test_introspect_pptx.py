"""Tests for introspect_pptx.py (User Story 3).

After the Phase 10 reverse-spec redesign the mock generator produces a
corpus rather than a single PPTX. These tests pick the Week 1 deck as a
representative single-PPTX target.
"""

from __future__ import annotations

import sys
import unittest
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_ROOT))
sys.path.insert(0, str(REPO_ROOT / "scripts"))

import introspect_pptx  # noqa: E402
from tests import conftest_helpers  # noqa: E402


TMP = REPO_ROOT / "tests" / "_tmp"


class IntrospectTests(unittest.TestCase):

    @classmethod
    def setUpClass(cls) -> None:
        TMP.mkdir(parents=True, exist_ok=True)
        corpus = conftest_helpers.make_mock_corpus(TMP / "introspect_mock")
        cls.pptx = conftest_helpers.first_pptx(corpus)
        cls.default_path = TMP / "introspect_report_default.txt"
        cls.verbose_path = TMP / "introspect_report_verbose.txt"
        rc = introspect_pptx.main(["--input", str(cls.pptx), "--out", str(cls.default_path)])
        assert rc == 0
        cls.default_report = cls.default_path.read_text(encoding="utf-8")
        rc = introspect_pptx.main(
            ["--input", str(cls.pptx), "--out", str(cls.verbose_path), "--verbose"]
        )
        assert rc == 0
        cls.verbose_report = cls.verbose_path.read_text(encoding="utf-8")

    def test_summary_lists_extensions_seen_in_corpus(self) -> None:
        # The Week deck links to .glc, .docx/.png (analysis), and rarely .wav.
        self.assertIn(".glc:", self.default_report)
        # At least one of the analysis-sheet extensions must appear.
        self.assertTrue(".png:" in self.default_report or ".docx:" in self.default_report,
                        "Expected analysis-sheet extension in summary")
        # New shape-hyperlink split (live vs vestigial) appears in summary.
        self.assertIn("Shape-level hyperlinks (live):", self.default_report)
        self.assertIn("Shape-level hyperlinks (vestigial absolute file:///):",
                      self.default_report)
        self.assertIn("Text-run hyperlinks:", self.default_report)
        # Gram-aware counts.
        self.assertIn("Total grams extracted:", self.default_report)

    def test_default_per_gram_section_lists_grams(self) -> None:
        section_2 = self.default_report.split("=== Section 2: Per-gram ===", 1)[1]
        # No verbose sections in default mode.
        self.assertNotIn("=== Section 3:", section_2)
        self.assertIn("analysis:", section_2)
        self.assertIn("glc[", section_2)

    def test_default_mode_omits_verbose_sections(self) -> None:
        self.assertNotIn("=== Section 3: Per-slide", self.default_report)
        self.assertNotIn("=== Section 4: Hyperlink targets", self.default_report)

    def test_verbose_per_shape_section_records_position_and_text(self) -> None:
        section_3 = self.verbose_report.split("=== Section 3: Per-slide (verbose) ===", 1)[1]
        section_3 = section_3.split("=== Section 4:", 1)[0]
        self.assertIn("name=", section_3)
        self.assertIn("type=", section_3)
        self.assertIn("pos=(", section_3)
        self.assertIn("hyperlink=", section_3)

    def test_verbose_hyperlink_targets_section_groups_by_extension(self) -> None:
        section_4 = self.verbose_report.split(
            "=== Section 4: Hyperlink targets (verbose, raw) ===", 1)[1]
        self.assertIn("-- .glc --", section_4)
        # At least one of the two analysis-sheet extensions must be present.
        self.assertTrue("-- .png --" in section_4 or "-- .docx --" in section_4)

    def test_slides_filter_restricts_per_gram_section(self) -> None:
        path = TMP / "introspect_filtered.txt"
        rc = introspect_pptx.main([
            "--input", str(self.pptx),
            "--out", str(path),
            "--slides", "1",
        ])
        self.assertEqual(rc, 0)
        section_2 = path.read_text(encoding="utf-8").split(
            "=== Section 2: Per-gram ===", 1)[1]
        self.assertIn("Slide 1", section_2)
        # If the deck has more than one slide, slide 2 should be filtered out.
        self.assertNotIn("Slide 2", section_2)

    def test_slides_filter_also_restricts_verbose_sections(self) -> None:
        path = TMP / "introspect_filtered_verbose.txt"
        rc = introspect_pptx.main([
            "--input", str(self.pptx),
            "--out", str(path),
            "--slides", "1",
            "--verbose",
        ])
        self.assertEqual(rc, 0)
        section_3 = path.read_text(encoding="utf-8").split(
            "=== Section 3: Per-slide (verbose) ===", 1)[1].split(
            "=== Section 4:", 1)[0]
        self.assertIn("Slide 1", section_3)
        self.assertNotIn("Slide 2", section_3)

    def test_unexpected_shape_count_is_flagged(self) -> None:
        from pptx import Presentation
        from unittest.mock import patch
        all_records = []
        prs = Presentation(self.pptx)
        for i, slide in enumerate(prs.slides, start=1):
            all_records.extend(introspect_pptx.collect_shape_records(slide, slide_number=i))
        with patch.object(introspect_pptx, "EXPECTED_SHAPES_PER_CONTENT_SLIDE", 999), \
             patch.object(introspect_pptx, "SHAPE_DEVIATION_TOLERANCE", 0):
            text = introspect_pptx.render_summary(self.pptx.name, len(list(prs.slides)), all_records)
        self.assertIn("Slides flagged", text)
        self.assertNotIn("Slides flagged (deviating shape count): none", text)


if __name__ == "__main__":
    unittest.main()

"""Tests for mock_pptx.py — corpus-aware (Phase 10 redesign).

Asserts properties of the multi-publication corpus described in
``source/notes/reverse-spec.md`` §1, §3, §4, §7. The pre-reverse-spec
tests (T061–T065) that assumed a single PPTX with 15 grams per 3×5 slide
are gone; their assumptions no longer hold.
"""

from __future__ import annotations

import re
import sys
import unittest
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(REPO_ROOT))

import mock_pptx  # noqa: E402
from tests import conftest_helpers  # noqa: E402

from pptx import Presentation
from pptx.oxml.ns import qn

TMP = REPO_ROOT / "tests" / "_tmp"


def _shape_hyperlink(shape) -> str | None:
    nv_sp_pr = shape._element.find(qn("p:nvSpPr"))
    if nv_sp_pr is None:
        return None
    c_nv_pr = nv_sp_pr.find(qn("p:cNvPr"))
    if c_nv_pr is None:
        return None
    hlink = c_nv_pr.find(qn("a:hlinkClick"))
    if hlink is None:
        return None
    rel_id = hlink.get(qn("r:id"))
    if not rel_id:
        return None
    return shape.part.rels[rel_id].target_ref


def _run_hyperlinks(shape) -> list[str]:
    hrefs: list[str] = []
    if not shape.has_text_frame:
        return hrefs
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            try:
                if run.hyperlink.address:
                    hrefs.append(run.hyperlink.address)
            except Exception:
                pass
    return hrefs


class CorpusShapeTests(unittest.TestCase):

    @classmethod
    def setUpClass(cls) -> None:
        TMP.mkdir(parents=True, exist_ok=True)
        cls.corpus = conftest_helpers.make_mock_corpus(TMP / "mock_corpus")

    def test_one_pptx_per_publication(self) -> None:
        for pub in mock_pptx.PUBLICATIONS:
            expected = self.corpus / pub.name / f"{pub.name}.pptx"
            self.assertTrue(expected.is_file(), f"missing {expected}")

    def test_files_folder_contains_gram_subfolders(self) -> None:
        # Sample a non-batched publication.
        pub = mock_pptx.PUBLICATIONS[0]  # Week 1
        files_dir = self.corpus / pub.name / f"{pub.name} Files"
        self.assertTrue(files_dir.is_dir())
        gram_dirs = [p for p in files_dir.iterdir() if p.is_dir()]
        # Family target ±15% to accommodate gap simulation.
        target = pub.family.grams
        self.assertGreaterEqual(len(gram_dirs), int(target * 0.85))
        self.assertLessEqual(len(gram_dirs), int(target * 1.15) + 2)

    def test_pub10_uses_ten_gram_batched_folders(self) -> None:
        pub10 = next(p for p in mock_pptx.PUBLICATIONS if p.batched)
        files_dir = self.corpus / pub10.name / f"{pub10.name} Files"
        batches = sorted(p.name for p in files_dir.iterdir() if p.is_dir())
        self.assertTrue(batches, "Pub10 should have batch subfolders")
        for batch in batches:
            self.assertRegex(batch, r"^Pub 10_Ed 2_\(\d+-\d+\)$",
                             f"batch folder {batch!r} does not match expected pattern")
        # Grams live one level deeper inside batches.
        any_gram = next(
            (p for batch in files_dir.iterdir() if batch.is_dir()
             for p in batch.iterdir() if p.is_dir()),
            None,
        )
        self.assertIsNotNone(any_gram)

    def test_no_fr_variant_drops_fr_prefix(self) -> None:
        pub = next(p for p in mock_pptx.PUBLICATIONS if p.no_fr)
        prs = Presentation(self.corpus / pub.name / f"{pub.name}.pptx")
        descriptors = []
        for slide in prs.slides:
            for shape in slide.shapes:
                href = _shape_hyperlink(shape)
                if href and shape.has_text_frame:
                    text = "".join(r.text or "" for p in shape.text_frame.paragraphs for r in p.runs)
                    descriptors.append(text)
        self.assertTrue(descriptors, "No FR publication should have gram tiles")
        for d in descriptors:
            self.assertNotIn(": FR ", d, f"No-FR variant still has FR prefix in: {d!r}")

    def test_gram_tile_descriptor_uses_colon_split(self) -> None:
        # Reverse-spec §4 — every gram rectangle text matches "Gram N: ...".
        pub = mock_pptx.PUBLICATIONS[0]
        prs = Presentation(self.corpus / pub.name / f"{pub.name}.pptx")
        found = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if _shape_hyperlink(shape) and shape.has_text_frame:
                    text = "".join(r.text or "" for p in shape.text_frame.paragraphs for r in p.runs).strip()
                    self.assertRegex(text, r"^Gram \d+: ",
                                     f"gram tile descriptor missing colon-split: {text!r}")
                    found += 1
        self.assertGreater(found, 0)

    def test_lofar_count_per_gram_in_1_to_4_range(self) -> None:
        pub = mock_pptx.PUBLICATIONS[0]
        prs = Presentation(self.corpus / pub.name / f"{pub.name}.pptx")
        for slide in prs.slides:
            # Count run hyperlinks per non-header text frame.
            for shape in slide.shapes:
                if _shape_hyperlink(shape):
                    continue
                hrefs = _run_hyperlinks(shape)
                if not hrefs:
                    continue
                self.assertGreaterEqual(len(hrefs), 1)
                self.assertLessEqual(len(hrefs), 4)
                for h in hrefs:
                    self.assertTrue(h.endswith(".glc") or h.endswith(".wav"),
                                    f"unexpected Lofar target: {h}")

    def test_analysis_sheet_mix_includes_both_docx_and_png(self) -> None:
        # Aggregated across the whole corpus.
        kinds = {"docx": 0, "png": 0}
        for pub in mock_pptx.PUBLICATIONS:
            prs = Presentation(self.corpus / pub.name / f"{pub.name}.pptx")
            for slide in prs.slides:
                for shape in slide.shapes:
                    href = _shape_hyperlink(shape)
                    if href is None:
                        continue
                    if href.lower().endswith(".docx"):
                        kinds["docx"] += 1
                    elif href.lower().endswith(".png"):
                        kinds["png"] += 1
        self.assertGreater(kinds["docx"], 0)
        self.assertGreater(kinds["png"], 0)
        # Aim is roughly 50/50; allow a wide tolerance.
        ratio = kinds["docx"] / (kinds["docx"] + kinds["png"])
        self.assertGreater(ratio, 0.3)
        self.assertLess(ratio, 0.7)

    def test_welcome_and_exit_slides_bracket_each_publication(self) -> None:
        for pub in mock_pptx.PUBLICATIONS:
            prs = Presentation(self.corpus / pub.name / f"{pub.name}.pptx")
            slides = list(prs.slides)
            self.assertGreaterEqual(len(slides), 3,
                                    f"{pub.name} should have welcome + content + exit")

            welcome_text = "".join(
                r.text or "" for s in slides[0].shapes if s.has_text_frame
                for p in s.text_frame.paragraphs for r in p.runs
            )
            self.assertIn(f"Welcome to {pub.name}", welcome_text)
            self.assertIn("Instructor Version", welcome_text)
            # Welcome slide carries no shape-level hyperlinks (no gram content).
            self.assertFalse(
                any(_shape_hyperlink(s) for s in slides[0].shapes),
                f"{pub.name} welcome slide should have no hyperlinks",
            )

            exit_text = "".join(
                r.text or "" for s in slides[-1].shapes if s.has_text_frame
                for p in s.text_frame.paragraphs for r in p.runs
            )
            self.assertIn(f"End of {pub.name}", exit_text)
            self.assertIn("Instructor Version", exit_text)
            self.assertFalse(
                any(_shape_hyperlink(s) for s in slides[-1].shapes),
                f"{pub.name} exit slide should have no hyperlinks",
            )

    def test_gram_numbering_has_gaps(self) -> None:
        # At least one publication should show a non-contiguous gram-number sequence.
        any_gaps = False
        for pub in mock_pptx.PUBLICATIONS:
            files_dir = self.corpus / pub.name / f"{pub.name} Files"
            if pub.batched:
                continue  # batched naming makes gap detection awkward
            nums = sorted(int(re.search(r"\d+", p.name).group())
                          for p in files_dir.iterdir() if p.is_dir())
            if nums and (max(nums) - min(nums) + 1) > len(nums):
                any_gaps = True
                break
        self.assertTrue(any_gaps, "Expected at least one publication with gram-number gaps")


if __name__ == "__main__":
    unittest.main()

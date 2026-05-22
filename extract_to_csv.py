"""Extractor (User Story 2): walk a content tree and emit the intermediate CSV.

Stage 2 of the migration pipeline. Walks every ``.pptx`` under the
configured input root, classifies each as ``main`` or
``progress-test-N`` by filename pattern (R2), opens each via
``python-pptx``, and writes one CSV row per resulting DITA topic
(R11, contracts/csv-schema.md). The shape-grouping step itself is the
documented ``NotImplementedError`` stub mandated by FR-015 and R1; every
piece of surrounding infrastructure is fully implemented so the rest of
the team can write tests against it before the stub is replaced.

Logging follows R10: dual stdout + ``extract.log``, three levels
(INFO/WARNING/ERROR), no silent exception swallowing (FR-014).
"""

from __future__ import annotations

import argparse
import csv
import logging
import re
import sys
import urllib.parse
import xml.etree.ElementTree as ET
from collections import Counter, defaultdict
from dataclasses import dataclass, field
from pathlib import Path, PurePosixPath, PureWindowsPath
from typing import Iterable, Iterator

from pptx import Presentation
from pptx.oxml.ns import qn


CSV_COLUMNS: tuple[str, ...] = (
    "publication", "chapter", "gram_id", "vessel_name", "topic_type",
    "sequence", "topic_filename", "display_text", "link_href", "glc_path",
    "time_end", "freq_end", "png_path", "wav_treatment", "warnings",
)

DEFAULT_TEST_PATTERN: str = "progress test"

# Prefixes that identify the welcome / exit framing slides emitted by
# ``mock_pptx.py``. These slides carry no gram content and must not
# contribute rows to the CSV.
FRAMING_TITLE_PREFIXES: tuple[str, ...] = ("Welcome to ", "End of ")

# A gram header's shape-level hyperlink must target one of these
# extensions. The audited legacy corpora use ``analysis sheet.doc`` or
# the newer ``*ANALYSIS.png`` variant; ``.docx`` is included for forward
# compatibility with re-authored decks. A ``.glc`` target on a
# shape-level link is treated as an authoring residue, not a header.
ANALYSIS_SHEET_EXTENSIONS: tuple[str, ...] = (".doc", ".docx", ".png")

LOGGER = logging.getLogger(__name__)


# -----------------------------------------------------------------------------
# Logging convention (R10) -- mirrored in generate_dita.py.
# -----------------------------------------------------------------------------

def setup_logging(log_path: Path) -> None:
    """Configure dual stdout + per-stage-file logging."""
    root = logging.getLogger()
    for handler in list(root.handlers):
        root.removeHandler(handler)
    formatter = logging.Formatter("%(asctime)s %(levelname)s %(name)s: %(message)s")
    stream = logging.StreamHandler(sys.stdout)
    stream.setLevel(logging.INFO)
    stream.setFormatter(formatter)
    file_handler = logging.FileHandler(log_path, mode="w", encoding="utf-8")
    file_handler.setLevel(logging.DEBUG)
    file_handler.setFormatter(formatter)
    root.setLevel(logging.DEBUG)
    root.addHandler(stream)
    root.addHandler(file_handler)


# -----------------------------------------------------------------------------
# GLC parser (R6, contracts/glc-schema.md)
# -----------------------------------------------------------------------------

@dataclass
class GlcDocument:
    image_filename: str = ""
    time_end: str = ""
    freq_end: str = ""
    warnings: list[str] = field(default_factory=list)


def parse_glc(path: Path) -> GlcDocument:
    """Parse a GLC XML file tolerantly. Never raises.

    Per contracts/glc-schema.md: missing elements yield empty values plus
    a verbatim warning string; malformed XML yields an empty result with
    a single ``"GLC malformed: <reason>"`` warning. Path stripping uses
    ``pathlib.PureWindowsPath(raw).name`` so a Windows ``W:\\foo\\bar.PNG``
    surfaces as ``bar.PNG``.
    """
    doc = GlcDocument()
    try:
        tree = ET.parse(path)
    except ET.ParseError as exc:
        first = str(exc).splitlines()[0] if str(exc) else "unknown error"
        doc.warnings.append(f"GLC malformed: {first}")
        return doc
    except OSError as exc:
        doc.warnings.append(f"GLC malformed: {exc}")
        return doc

    root = tree.getroot()

    filenames = root.findall("data_source/filename")
    if not filenames:
        doc.warnings.append("GLC missing filename")
    else:
        if len(filenames) > 1:
            doc.warnings.append("GLC duplicate filename")
        raw = (filenames[0].text or "").strip()
        if not raw:
            doc.warnings.append("GLC missing filename")
            doc.image_filename = ""
        else:
            doc.image_filename = PureWindowsPath(raw).name

    bottom = root.findtext("data_source/bitmap_crop_values/bottom_crop")
    if bottom is None or not bottom.strip():
        doc.warnings.append("GLC missing bottom_crop")
        doc.time_end = ""
    else:
        doc.time_end = bottom.strip()

    bandwidth = root.findtext("settings/lofar/bandwidth")
    if bandwidth is None or not bandwidth.strip():
        doc.warnings.append("GLC missing bandwidth")
        doc.freq_end = ""
    else:
        doc.freq_end = bandwidth.strip()

    return doc


# -----------------------------------------------------------------------------
# Path resolution and walking
# -----------------------------------------------------------------------------

def resolve_glc_path(href: str, content_root: Path, source_dir: Path | None = None) -> Path | None:
    """Resolve a GLC href against the per-gram or per-ten-grams layout (FR-006).

    Returns ``None`` and logs a WARNING when the file cannot be found.
    """
    if not href:
        LOGGER.warning("GLC not found: empty href")
        return None
    rel = Path(href.replace("\\", "/"))
    candidates: list[Path] = []
    if source_dir is not None:
        candidates.append((source_dir / rel).resolve(strict=False))
    candidates.append((content_root / rel).resolve(strict=False))
    candidates.append((content_root / rel.name).resolve(strict=False))
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    LOGGER.warning("GLC not found: %s", href)
    return None


def _rel_to_root(path: Path, content_root: Path) -> str:
    """Return ``path`` as a POSIX string relative to ``content_root`` when possible."""
    root_abs = content_root.resolve()
    if path.is_relative_to(root_abs):
        return path.relative_to(root_abs).as_posix()
    return path.as_posix()


def resolve_asset_path(href: str, content_root: Path, source_dir: Path | None) -> str:
    """Resolve an asset href (PNG/DOCX/WAV) to a path relative to ``content_root``.

    Returns the resolved POSIX path string when the file is found, or the
    original href stripped of backslashes when it is not — the generator
    treats a missing asset as a dangling reference, not a fatal error.
    """
    if not href:
        return ""
    resolved = resolve_glc_path(href, content_root, source_dir=source_dir)
    if resolved is not None:
        return _rel_to_root(resolved, content_root)
    return href.replace("\\", "/")


def walk_pptxs(input_root: Path) -> Iterator[Path]:
    """Yield every ``.pptx`` under ``input_root`` in deterministic sorted order."""
    yield from sorted(input_root.rglob("*.pptx"))


# -----------------------------------------------------------------------------
# Slug + chapter helpers (R3)
# -----------------------------------------------------------------------------

_SLUG_NON_ALNUM = re.compile(r"[^a-z0-9]+")


def slugify(text: str) -> str:
    """Lower-case, ASCII-only, hyphen-separated slug with collapsed runs."""
    ascii_only = text.encode("ascii", "ignore").decode("ascii").lower()
    slug = _SLUG_NON_ALNUM.sub("-", ascii_only).strip("-")
    return slug


def classify_publication(
    pptx: Path,
    test_pattern: str,
    allocated: dict[str, int],
) -> tuple[str, str | None, str | None]:
    """Return ``(publication, chapter, chapter_slug)`` per R2/R3.

    Progress-test PPTXs are detected by case-insensitive substring match
    against the filename. Test numbering is allocated stably in the order
    callers request previously-unseen test PPTXs.
    """
    name = pptx.name.lower()
    if test_pattern.lower() in name:
        if pptx.stem not in allocated:
            allocated[pptx.stem] = len(allocated) + 1
        return (f"progress-test-{allocated[pptx.stem]}", None, None)
    chapter_title = pptx.parent.name
    return ("main", chapter_title, slugify(chapter_title))


# -----------------------------------------------------------------------------
# Shape grouping stub (FR-015 / R1)
# -----------------------------------------------------------------------------

@dataclass
class GlcLink:
    display_text: str
    href: str


@dataclass
class GramPlaceholder:
    gram_id: str
    vessel_name: str
    png_href: str | None
    glc_links: list[GlcLink]


def _shape_level_hyperlink(shape) -> str | None:
    """Return the shape-level hyperlink target, or None.

    Walks the lxml element directly so we don't depend on python-pptx
    exposing a high-level accessor (it doesn't, for shape-level clicks).
    Searches every descendant ``p:cNvPr`` so the lookup works regardless
    of the shape's outer XML wrapper — ``p:sp`` (autoshape/textbox),
    ``p:pic`` (picture), ``p:cxnSp`` (connector), or
    ``p:graphicFrame`` (chart/table/etc.). The actual ``a:hlinkClick``
    always sits under ``cNvPr`` regardless of wrapper.
    """
    for c_nv_pr in shape._element.iter(qn("p:cNvPr")):
        hlink = c_nv_pr.find(qn("a:hlinkClick"))
        if hlink is None:
            continue
        rel_id = hlink.get(qn("r:id"))
        if not rel_id:
            continue
        try:
            return shape.part.rels[rel_id].target_ref
        except KeyError:
            continue
    return None


def _run_hyperlinks_in_shape(shape) -> list[tuple[str, str]]:
    """Return a list of ``(visible_text, href)`` per hyperlinked text run.

    Only runs whose ``a:hlinkClick`` resolves to an external relationship
    target are returned. Plain runs without hyperlinks are skipped.
    """
    pairs: list[tuple[str, str]] = []
    if not shape.has_text_frame:
        return pairs
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                href = run.hyperlink.address
            except Exception:
                href = None
            if href:
                pairs.append((run.text or "", href))
    return pairs


def _gram_folder_key(href: str) -> str:
    """Last directory segment of a hyperlink target, URL-decoded and
    lowercased. Used to associate each .glc with the gram-header that
    points at the same ``GramNN/`` directory on disk.

    Returns the empty string if the href has no parent directory
    (e.g. a bare filename) — such links can't be folder-matched and
    are skipped by the caller.
    """
    if not href:
        return ""
    decoded = urllib.parse.unquote(href.split("?", 1)[0].split("#", 1)[0])
    path = PurePosixPath(decoded.replace("\\", "/"))
    parts = path.parts
    if len(parts) < 2:
        return ""
    return parts[-2].lower()


def _bbox(shape) -> tuple[int, int, int, int]:
    """Return ``(left, top, right, bottom)`` in EMUs."""
    left = shape.left or 0
    top = shape.top or 0
    width = shape.width or 0
    height = shape.height or 0
    return (left, top, left + width, top + height)


def _iter_leaf_shapes(shapes):
    """Yield every leaf shape on a slide, expanding GROUPs (shape_type=6)."""
    for shape in shapes:
        if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def is_framing_slide(slide) -> bool:
    """Return True for welcome / exit slides that should be skipped.

    Detection is by title text prefix (``"Welcome to "`` / ``"End of "``)
    so it survives slide-position shuffles and is robust against real
    decks that may add or remove framing slides.
    """
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text_frame.text or "").strip()
        if not text:
            continue
        if text.startswith(FRAMING_TITLE_PREFIXES):
            return True
    return False


def extract_grams_from_slide(slide, slide_num: int) -> list[GramPlaceholder]:
    """Return the gram placeholders on ``slide`` per reverse-spec §4.

    Each gram tile has:

    - a header shape carrying a *shape-level* hyperlink to an Analysis
      Sheet (``.docx`` or ``.png``); its text is a ``"Gram N: <detail>"``
      descriptor (split at the first colon — left side is the
      student-visible label, right side is instructor-visible);
    - one or more text-frame shapes positioned beneath the header
      containing text *runs* hyperlinked to ``.glc`` files in the same
      gram folder. Every Lofar text-run hyperlink in the audited corpus
      targets a ``.glc``; the ``.wav`` case is one indirection deeper,
      inside the ``.glc``'s ``data_source/filename`` element (see
      ``high-level-spec.md`` §1.5).

    Grouping is done by spatial proximity: for each header we collect
    every Lofar text run carried by text-frame shapes whose bounding
    box sits below the header and overlaps its horizontal extent. This
    survives renamed shapes, ad-hoc shape ordering, and GROUP wrappers.
    """
    leaves = list(_iter_leaf_shapes(slide.shapes))

    # 1) Identify headers (shape-level hyperlink to an analysis-sheet
    #    asset, on a text-bearing shape).
    headers: list[tuple[object, str]] = []  # (shape, analysis_href)
    for shape in leaves:
        href = _shape_level_hyperlink(shape)
        if href is None:
            continue
        if not shape.has_text_frame:
            continue
        # Skip vestigial overlay shapes whose shape-level hyperlinks point
        # at absolute file:/// URIs from a long-gone authoring path. The
        # live header buttons use paths relative to the PPTX; the dead
        # overlay (e.g. Group 197/Rectangle children) is left over from
        # earlier authoring iterations and never resolves.
        if href.lower().startswith("file:///"):
            continue
        # Whitelist the analysis-sheet extensions. Legacy authoring
        # sometimes promotes a .glc text-run hyperlink to the shape
        # level on the small .glc-bearing shapes; without this filter
        # those would be mistaken for gram headers.
        if not href.lower().endswith(ANALYSIS_SHEET_EXTENSIONS):
            continue
        headers.append((shape, href))

    # 2) Identify candidate Lofar shapes — shapes whose hyperlink targets
    #    a .glc. Two authoring styles both occur in the legacy corpus:
    #      (a) text-bearing shapes (autoshape/textbox) with the .glc
    #          hyperlink on a text *run* — the common case;
    #      (b) picture shapes with the .glc hyperlink at the shape level
    #          — used for at least one gram per slide in some decks
    #          (e.g. the gram whose analysis sheet is "V III .doc").
    candidates: list[tuple[object, list[tuple[str, str]]]] = []
    header_ids = {id(s) for s, _ in headers}
    for shape in leaves:
        if id(shape) in header_ids:
            continue
        keep: list[tuple[str, str]] = []
        # (a) text-run .glc hyperlinks
        for text, href in _run_hyperlinks_in_shape(shape):
            if href.lower().endswith(".glc"):
                keep.append((text, href))
            else:
                LOGGER.warning(
                    "Lofar text run hyperlinks to a non-.glc target (%r); "
                    "expected .glc — row dropped", href,
                )
        # (b) shape-level .glc hyperlink (picture-style). Skip vestigial
        # absolute file:/// URIs the same way the header step does.
        shape_href = _shape_level_hyperlink(shape)
        if (
            shape_href
            and shape_href.lower().endswith(".glc")
            and not shape_href.lower().startswith("file:///")
        ):
            display = ""
            if getattr(shape, "has_text_frame", False):
                display = (shape.text_frame.text or "").strip()
            keep.append((display, shape_href))
        if keep:
            candidates.append((shape, keep))

    # 3) Match each .glc to its gram header by shared parent folder on
    #    disk (e.g. both targets live under ``.../Gram001/``). This is
    #    far more robust than spatial proximity: it survives shapes
    #    moved off-screen, off-grid layouts, hidden overlay rectangles,
    #    and decks where the .glc shapes don't sit directly beneath
    #    their header rectangle. Per the corpus survey, every live
    #    header and its .glc children share a single ``GramNN/`` parent.
    folder_to_header: dict[str, tuple[object, str]] = {}
    for header, href in headers:
        key = _gram_folder_key(href)
        if not key:
            h_left, h_top, _, _ = _bbox(header)
            LOGGER.warning(
                "Slide %d: header at top=%d left=%d has no gram-folder in href %r",
                slide_num, h_top, h_left, href,
            )
            continue
        if key in folder_to_header:
            LOGGER.warning(
                "Slide %d: duplicate gram folder %r — second header for %r ignored",
                slide_num, key, href,
            )
            continue
        folder_to_header[key] = (header, href)

    header_to_pairs: dict[int, list[tuple[str, str]]] = defaultdict(list)
    for cand, pairs in candidates:
        for text, glc_href in pairs:
            key = _gram_folder_key(glc_href)
            if not key:
                LOGGER.warning(
                    "Slide %d: .glc %r has no gram-folder in href",
                    slide_num, glc_href,
                )
                continue
            hdr_entry = folder_to_header.get(key)
            if hdr_entry is None:
                LOGGER.warning(
                    "Slide %d: .glc %r names folder %r but no matching header on this slide",
                    slide_num, glc_href, key,
                )
                continue
            header_id = id(hdr_entry[0])
            header_to_pairs[header_id].append((text, glc_href))

    # 4) Build grams in reading order (top, left) for stable CSV output.
    grams: list[GramPlaceholder] = []
    headers.sort(key=lambda hb: (hb[0].top or 0, hb[0].left or 0))

    for header, analysis_href in headers:
        lofar_pairs = header_to_pairs.get(id(header), [])
        if not lofar_pairs:
            h_left, h_top, _, _ = _bbox(header)
            LOGGER.warning(
                "Slide %d: gram header at top=%d left=%d has no Lofar box",
                slide_num, h_top, h_left,
            )

        descriptor = "".join(
            run.text or "" for para in header.text_frame.paragraphs for run in para.runs
        ).strip()
        gram_id, instructor_detail = _split_descriptor(descriptor)

        glc_links = [GlcLink(display_text=t.strip(), href=h) for t, h in lofar_pairs]
        grams.append(GramPlaceholder(
            gram_id=gram_id,
            vessel_name=instructor_detail,
            png_href=analysis_href,
            glc_links=glc_links,
        ))
    return grams


def _split_descriptor(descriptor: str) -> tuple[str, str]:
    """Split a `"Gram N: <detail>"` descriptor at the first colon.

    Returns ``(gram_id, instructor_detail)``. If no colon is present, the
    whole text is treated as the gram_id and instructor_detail is empty.
    Gram_id is normalised to a plain integer string (``"7"``, ``"12"``)
    when the left side parses as ``"Gram N"`` — the canonical CSV form
    per csv-schema.md so authors can renumber with a bare number when
    refactoring content between chapters.
    """
    # Collapse runs of whitespace (legacy decks pad gram titles with
    # multi-space sequences to force in-shape line breaks).
    descriptor = " ".join(descriptor.split())
    if not descriptor:
        return ("", "")
    left, sep, right = descriptor.partition(":")
    left = left.strip()
    right = right.strip() if sep else ""
    m = re.match(r"^Gram\s+(\d+)$", left, re.IGNORECASE)
    if m:
        left = str(int(m.group(1)))
    return (left, right)


# -----------------------------------------------------------------------------
# Row construction
# -----------------------------------------------------------------------------

def _gram_num_from_id(gram_id: str) -> str:
    digits = re.findall(r"\d+", gram_id)
    return digits[0] if digits else "00"


def gram_to_rows(
    gram: GramPlaceholder,
    publication: str,
    chapter: str | None,
    chapter_slug: str | None,
    content_root: Path,
    source_dir: Path,
) -> list[dict]:
    """Expand one gram into N+1 CSV rows (N GLC links + 1 analysis row)."""
    rows: list[dict] = []
    gram_num = _gram_num_from_id(gram.gram_id)

    for i, link in enumerate(gram.glc_links, start=1):
        warnings: list[str] = []
        href = link.href
        glc_path = ""
        time_end = ""
        freq_end = ""
        png_path = ""
        display_text = link.display_text
        resolved = resolve_glc_path(href, content_root, source_dir=source_dir)
        if resolved is None:
            warnings.append("GLC not found")
            glc_path = href
        else:
            glc_path = _rel_to_root(resolved, content_root)
            glc = parse_glc(resolved)
            warnings.extend(glc.warnings)
            time_end = glc.time_end
            freq_end = glc.freq_end
            if glc.image_filename:
                # The inner asset (.png/.jpg pre-rendered spectrogram, or
                # .wav for live render) sits next to the GLC on disk.
                # Resolve against the GLC's directory so the path is
                # image-root-relative and the generator can copy it
                # directly (see dita-topic-schema.md §1.2/§1.3 for how
                # the generator dispatches on the extension).
                png_path = resolve_asset_path(
                    glc.image_filename, content_root, source_dir=resolved.parent,
                )

        rows.append({
            "publication": publication,
            "chapter": chapter or "",
            "gram_id": gram.gram_id,
            "vessel_name": gram.vessel_name,
            "topic_type": "glc",
            "sequence": str(i),
            "topic_filename": f"gram_{gram_num}.dita",
            "display_text": display_text,
            "link_href": href,
            "glc_path": glc_path,
            "time_end": time_end,
            "freq_end": freq_end,
            "png_path": png_path,
            "wav_treatment": "",
            "warnings": ", ".join(warnings),
        })

    analysis_warnings: list[str] = []
    analysis_png = gram.png_href or ""
    if not analysis_png:
        analysis_warnings.append("missing analysis PNG hyperlink")
        analysis_png_resolved = ""
    else:
        analysis_png_resolved = resolve_asset_path(analysis_png, content_root, source_dir)
    rows.append({
        "publication": publication,
        "chapter": chapter or "",
        "gram_id": gram.gram_id,
        "vessel_name": gram.vessel_name,
        "topic_type": "analysis",
        "sequence": "1",
        "topic_filename": f"gram_{gram_num}.dita",
        "display_text": "",
        "link_href": "",
        "glc_path": "",
        "time_end": "",
        "freq_end": "",
        "png_path": analysis_png_resolved,
        "wav_treatment": "",
        "warnings": ", ".join(analysis_warnings),
    })
    return rows


# -----------------------------------------------------------------------------
# CSV writer (R11)
# -----------------------------------------------------------------------------

def write_csv(rows: list[dict], out: Path) -> None:
    """Write rows to a UTF-8-with-BOM, CRLF, QUOTE_MINIMAL CSV."""
    out.parent.mkdir(parents=True, exist_ok=True)
    with out.open("w", encoding="utf-8-sig", newline="") as fh:
        writer = csv.DictWriter(
            fh, fieldnames=list(CSV_COLUMNS),
            quoting=csv.QUOTE_MINIMAL, lineterminator="\r\n",
        )
        writer.writeheader()
        for row in rows:
            writer.writerow({col: row.get(col, "") for col in CSV_COLUMNS})


# -----------------------------------------------------------------------------
# Orchestration
# -----------------------------------------------------------------------------

def main(argv: Iterable[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Extract PPTX content into the intermediate CSV")
    parser.add_argument("--input-root", required=True, type=Path, dest="input_root")
    parser.add_argument("--out", required=True, type=Path)
    parser.add_argument("--test-pattern", default=DEFAULT_TEST_PATTERN, dest="test_pattern")
    args = parser.parse_args(list(argv) if argv is not None else None)

    setup_logging(Path("extract.log"))

    if not args.input_root.is_dir():
        LOGGER.error("Input root does not exist or is not a directory: %s", args.input_root)
        return 1

    rows: list[dict] = []
    pptx_count = 0
    warning_counter: Counter[str] = Counter()
    allocated: dict[str, int] = {}

    try:
        for pptx in walk_pptxs(args.input_root):
            pptx_count += 1
            LOGGER.info("Processing PPTX %s", pptx)
            publication, chapter, chapter_slug = classify_publication(pptx, args.test_pattern, allocated)
            try:
                prs = Presentation(pptx)
            except Exception as exc:
                LOGGER.error("Cannot open PPTX %s: %s", pptx, exc)
                return 1
            for slide_num, slide in enumerate(prs.slides, start=1):
                if is_framing_slide(slide):
                    LOGGER.info("Skipping framing slide %d in %s", slide_num, pptx.name)
                    continue
                grams = extract_grams_from_slide(slide, slide_num)
                for gram in grams:
                    gram_rows = gram_to_rows(
                        gram, publication, chapter, chapter_slug,
                        args.input_root, source_dir=pptx.parent,
                    )
                    rows.extend(gram_rows)
                    for r in gram_rows:
                        if r["warnings"]:
                            for w in r["warnings"].split(", "):
                                warning_counter[w] += 1
    except NotImplementedError as exc:
        LOGGER.error("Shape grouping stub reached: %s", exc)
        # Still write whatever rows were collected before the stub fired.
        write_csv(rows, args.out)
        return 1

    write_csv(rows, args.out)

    distinct = ", ".join(f"{w}={c}" for w, c in sorted(warning_counter.items()))
    LOGGER.info(
        "Extraction summary: pptx=%d rows=%d warnings=%d distinct=[%s]",
        pptx_count, len(rows), sum(warning_counter.values()), distinct,
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())

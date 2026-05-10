"""Synthetic instructor PPTX generator (User Story 4).

Builds a deck shaped like a real AAAC training instructor presentation so
that the introspection, extraction, and DITA generator tests can run
without a binary fixture in the repository. The generator covers every
structural case the source spec mandates: a welcome slide; content slides
with a 3 row x 5 col grid of gram placeholders; vessel-named title
rectangles carrying shape-level click actions to analysis PNGs; link text
boxes carrying text-run hyperlinks to GLC and (for the configured WAV
overrides) WAV targets; varying GLC link counts per gram.

Run as a script:

    python mock_pptx.py --out tests/_tmp/mock.pptx

The constants below (vessel names, link-count tiers, WAV grams, slide and
grid geometry) are exported so the test suite can assert against them
without magic-number drift between mock and tests (R5).
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree


# -----------------------------------------------------------------------------
# Module-level constants (R5)
# -----------------------------------------------------------------------------

VESSEL_NAMES: tuple[str, ...] = (
    "Nordik Jockey", "Arctic Surveyor", "Baltic Trader", "Cape Hopper",
    "Drift Wanderer", "Eastern Voyager", "Frosted Fjord", "Glacier Hauler",
    "Harbour Master", "Iceberg Runner", "Jutland Express", "Kelp Forest",
    "Lighthouse Keeper", "Maritime Crown", "Northern Star", "Ocean Pilot",
    "Polar Companion", "Quayside Belle", "Reef Patroller", "Saltwater Lark",
    "Tidal Mariner", "Undertow Drifter", "Vessel Aurora", "Whaler Echo",
    "Xenia Foam", "Yarrow Sound", "Zephyr Tide", "Anchor Bay",
    "Beacon Shoal", "Compass Reach",
)

# Gram numbering 1..30 in three tiers with different link-count variation.
# Tier (low, high) is inclusive on both ends.
LINK_COUNT_BY_GRAM_RANGE: tuple[tuple[int, int, int], ...] = (
    (1, 10, 1),    # grams 1..10 -> 1 GLC link
    (11, 25, 2),   # grams 11..25 -> 2 GLC links
    (26, 30, 4),   # grams 26..30 -> 4 GLC links
)

WAV_GRAMS: tuple[int, ...] = (5, 20)

TOTAL_GRAMS: int = 30
GRAMS_PER_SLIDE: int = 15
GRID_ROWS: int = 3
GRID_COLS: int = 5
SLIDE_WIDTH_IN: float = 13.33
SLIDE_HEIGHT_IN: float = 7.5
GRID_TOP_MARGIN_IN: float = 1.0
GRID_LEFT_MARGIN_IN: float = 0.4
TITLE_HEIGHT_IN: float = 0.5
LINK_BOX_HEIGHT_IN: float = 1.1
CELL_GAP_IN: float = 0.1


# -----------------------------------------------------------------------------
# Derived counts (referenced by tests)
# -----------------------------------------------------------------------------

CONTENT_SLIDE_COUNT: int = (TOTAL_GRAMS + GRAMS_PER_SLIDE - 1) // GRAMS_PER_SLIDE
TOTAL_SLIDE_COUNT: int = 1 + CONTENT_SLIDE_COUNT


def link_count_for_gram(gram_num: int) -> int:
    """Return the number of GLC links the mock attaches to ``gram_num``."""
    for low, high, count in LINK_COUNT_BY_GRAM_RANGE:
        if low <= gram_num <= high:
            return count
    raise ValueError(f"Gram {gram_num} is outside the configured ranges")


# -----------------------------------------------------------------------------
# Hyperlink helpers
# -----------------------------------------------------------------------------

def add_shape_level_hyperlink(shape, target: str) -> None:
    """Attach a shape-level click action to ``shape`` pointing at ``target``.

    python-pptx exposes text-run hyperlinks but no high-level API for shape
    click actions, so the underlying lxml element gains an ``a:hlinkClick``
    under ``p:nvSpPr/p:nvPr`` and the relationship is registered on the
    slide part.
    """
    rel_id = shape.part.relate_to(
        target,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
        is_external=True,
    )
    nv_sp_pr = shape._element.find(qn("p:nvSpPr"))
    nv_pr = nv_sp_pr.find(qn("p:nvPr"))
    hlink = etree.SubElement(nv_pr, qn("a:hlinkClick"))
    hlink.set(qn("r:id"), rel_id)


def add_text_run_hyperlink(run, target: str) -> None:
    """Attach a text-run hyperlink to ``run`` pointing at ``target``."""
    run.hyperlink.address = target


# -----------------------------------------------------------------------------
# Slide builders
# -----------------------------------------------------------------------------

def add_welcome_slide(prs: Presentation) -> None:
    """Add the welcome slide (slide 1)."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title is not None:
        slide.shapes.title.text = "Welcome to AAAC Training Module 3"
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Instructor Version"
            break


def build_gram_placeholder(
    slide,
    row: int,
    col: int,
    gram_num: int,
    vessel: str,
    link_count: int,
    is_wav: bool,
) -> None:
    """Place one gram (title rectangle + link text box) on ``slide``.

    The title rectangle gets a shape-level click action to
    ``../images/gramNN_analysis.png``. The link text box gets one
    text-run hyperlink per requested link to either
    ``../gramNN/config_M.glc`` or (if ``is_wav``) a single WAV target.
    """
    cell_w = (SLIDE_WIDTH_IN - 2 * GRID_LEFT_MARGIN_IN) / GRID_COLS
    cell_h = (SLIDE_HEIGHT_IN - GRID_TOP_MARGIN_IN - 0.4) / GRID_ROWS
    left = Inches(GRID_LEFT_MARGIN_IN + col * cell_w)
    top = Inches(GRID_TOP_MARGIN_IN + row * cell_h)
    width = Inches(cell_w - CELL_GAP_IN)

    title_top = top
    title_height = Inches(TITLE_HEIGHT_IN)
    title_shape = slide.shapes.add_shape(1, left, title_top, width, title_height)
    title_shape.text_frame.text = f"Gram {gram_num:02d} - {vessel}"
    add_shape_level_hyperlink(title_shape, f"../images/gram{gram_num:02d}_analysis.png")

    link_top = Inches(GRID_TOP_MARGIN_IN + row * cell_h + TITLE_HEIGHT_IN + 0.05)
    link_height = Inches(LINK_BOX_HEIGHT_IN)
    link_box = slide.shapes.add_textbox(left, link_top, width, link_height)
    tf = link_box.text_frame
    tf.word_wrap = True
    for i in range(link_count):
        if i == 0:
            paragraph = tf.paragraphs[0]
        else:
            paragraph = tf.add_paragraph()
        run = paragraph.add_run()
        run.text = f"LOFAR {i + 1}"
        if is_wav and i == 0:
            target = f"../gram{gram_num:02d}/clip_{i + 1}.wav"
        else:
            target = f"../gram{gram_num:02d}/config_{i + 1}.glc"
        add_text_run_hyperlink(run, target)


def add_content_slide(prs: Presentation, slide_num: int, gram_start: int, gram_end: int) -> None:
    """Add a content slide hosting grams ``gram_start..gram_end`` inclusive."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    for offset, gram_num in enumerate(range(gram_start, gram_end + 1)):
        row = offset // GRID_COLS
        col = offset % GRID_COLS
        vessel = VESSEL_NAMES[(gram_num - 1) % len(VESSEL_NAMES)]
        link_count = link_count_for_gram(gram_num)
        is_wav = gram_num in WAV_GRAMS
        build_gram_placeholder(slide, row, col, gram_num, vessel, link_count, is_wav)


# -----------------------------------------------------------------------------
# Orchestration
# -----------------------------------------------------------------------------

def main(argv: Iterable[str] | None = None) -> int:
    """Build and write the mock PPTX. Returns process exit code."""
    parser = argparse.ArgumentParser(description="Generate a synthetic instructor PPTX")
    parser.add_argument("--out", required=True, type=Path, help="Output .pptx path")
    args = parser.parse_args(list(argv) if argv is not None else None)

    out_path: Path = args.out
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    add_welcome_slide(prs)

    gram = 1
    slide_num = 2
    while gram <= TOTAL_GRAMS:
        gram_end = min(gram + GRAMS_PER_SLIDE - 1, TOTAL_GRAMS)
        add_content_slide(prs, slide_num, gram, gram_end)
        gram = gram_end + 1
        slide_num += 1

    prs.save(out_path)
    print(f"Wrote {out_path} with {TOTAL_SLIDE_COUNT} slides ({CONTENT_SLIDE_COUNT} content + 1 welcome).")
    return 0


if __name__ == "__main__":
    sys.exit(main())

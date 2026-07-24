"""Extractor (User Story 2): walk a content tree and emit the intermediate CSV.

Stage 2 of the migration pipeline. Walks every ``.pptx`` under the
configured input root, classifies each as ``main`` or
``progress-test-N`` by filename pattern (R2), opens each via
``python-pptx``, and writes one CSV row per resulting DITA topic
(R11, contracts/csv-schema.md). The shape-grouping step itself is the
documented ``NotImplementedError`` stub mandated by FR-015 and R1; every
piece of surrounding infrastructure is fully implemented so the rest of
the team can write tests against it before the stub is replaced.

Logging follows R10: dual stdout + ``extract.log``, three levels
(INFO/WARNING/ERROR), no silent exception swallowing (FR-014).
"""

from __future__ import annotations

import argparse
import csv
import logging
import re
import sys
import urllib.parse
import xml.etree.ElementTree as ET
from collections import Counter, defaultdict
from dataclasses import dataclass, field
from pathlib import Path, PurePosixPath, PureWindowsPath
from typing import Iterable, Iterator

from pptx import Presentation
from pptx.oxml.ns import qn


CSV_COLUMNS: tuple[str, ...] = (
    "publication", "chapter", "target_doc", "target_chapter", "gram_id", "vessel_name", "topic_type",
    "sequence", "topic_filename", "display_text", "link_href", "glc_path",
    "time_end", "bandwidth", "bandcentre", "png_path", "target_ext", "file_size", "wav_treatment", "warnings",
)

DEFAULT_TEST_PATTERN: str = "progress test"
DEFAULT_FINAL_PATTERN: str = "final assessment"
DEFAULT_JOINING_PATTERN: str = "joining"

# GramFrame needs the full time + frequency coordinate system to render a gram,
# so an *image* GLC-backed gram (pre-rendered .png/.jpg embedded inline) requires
# these three "view" fields. ``time_end`` is the image's pixel height (issue
# #148, derived at extraction from the file on disk); ``bandwidth``/``bandcentre``
# come from the .glc. A .wav-backed gram does *not* need them: it is surfaced
# downstream as a plain link to its .glc (the on-PC GLC viewer renders it live,
# reading the .glc directly), so no GramFrame table is emitted and the view
# fields are never consumed for it -- a blank one is fine. We catch a blank on an
# image row here so the failure surfaces at extraction (issue #92) rather than
# late and cryptically in dedupe. ``--relaxed`` substitutes ``RELAXED_DEFAULT``
# so the rest of the toolchain can still be exercised against an incomplete corpus.
GLC_VIEW_FIELDS: tuple[str, ...] = ("time_end", "bandwidth", "bandcentre")
# The subset of the view fields that come from the GLC author and so are subject
# to the strict extract-time fail-fast (issue #92). ``time_end`` is excluded: it
# is derived from the image's pixel height (issue #148), so a blank one is a
# dangling-asset problem, not an author omission — see ``glc_view_problems``.
GLC_AUTHORED_VIEW_FIELDS: tuple[str, ...] = ("bandwidth", "bandcentre")
RELAXED_DEFAULT: str = "100"

# Warning stamped on any row whose ``png_path`` references an asset (image or
# ``.wav``) that is not present on disk. Such a row dangles silently through the
# generator and only surfaces as a DITA-OT load error at publish time; catching
# it here at the extraction boundary lets the operator track and triage the
# instances (fix the source file, or drop the row) before publishing.
ASSET_MISSING_WARNING: str = "asset file missing on disk"

# Inner-asset extensions that feed a GramFrame table and so require the view
# fields. Only the embedded-image extensions qualify: a .wav is surfaced as a
# link to its .glc, not rendered by GramFrame, so it is exempt (its blank view
# fields are legitimate, not a defect). A GLC row with no resolved asset
# (target_ext == "") dangles per the missing-asset rule and is not subject to
# the check either. Mirrors the generator's dispatch (generate_dita.py).
GRAMFRAME_GLC_EXTENSIONS: tuple[str, ...] = (
    ".png", ".jpg", ".jpeg", ".gif")

# A demon marker (issue #151): ingest writes a ``demon.glc`` (and ``demon-2.glc``,
# … for several demons) into a gram folder beside a demon image. It is *not*
# hyperlinked from the slide, so ``extract`` discovers it by scanning the gram's
# first resolved Lofar folder (see ``_DEMON_MARKER_RE`` / ``_find_demon_markers``).

# Prefixes that identify the welcome / exit framing slides emitted by
# ``mock_pptx.py``. These slides carry no gram content and must not
# contribute rows to the CSV.
FRAMING_TITLE_PREFIXES: tuple[str, ...] = ("Welcome to ", "End of ")

# A gram header's shape-level hyperlink must target one of these
# extensions. The audited legacy corpora use ``analysis sheet.doc``,
# the newer ``*ANALYSIS.png`` variant, or a JPG export; ``.docx`` is
# included for forward compatibility with re-authored decks. A ``.glc``
# target on a shape-level link is treated as an authoring residue,
# not a header.
ANALYSIS_SHEET_EXTENSIONS: tuple[str, ...] = (".doc", ".docx", ".png", ".jpg", ".jpeg")

# Tokens that name an analysis sheet on disk, used by the Lofar-folder fallback
# (``_recover_analysis_sheet``) when a gram's analysis hyperlink is stale. Kept
# in sync with snapshot_analysis_docs.ANALYSIS_NAME_TOKEN /
# ANALYSIS_NAME_MISSPELLINGS so the fallback recognises exactly the sheets the
# snapshotter does. Only image sheets are recovered here — a Word sheet awaiting
# its snapshot render is a separate flow, left to Feature 007.
ANALYSIS_NAME_TOKEN = "analysis"
ANALYSIS_NAME_MISSPELLINGS: tuple[str, ...] = ("analaysis",)
ANALYSIS_IMAGE_EXTENSIONS: tuple[str, ...] = (".png", ".jpg", ".jpeg")

# Some legacy decks carry navigation buttons labelled like "N Questions"
# that link to an image (a self-test slide), not a gram. Their shape-level
# hyperlink targets a ``.png`` so they would otherwise be mistaken for a
# gram header. Any link whose visible label ends in the word "questions"
# (case-insensitive; the leading number is irrelevant) is silently ignored.
QUESTIONS_LABEL_WORD = "questions"

LOGGER = logging.getLogger(__name__)


# -----------------------------------------------------------------------------
# Logging convention (R10) -- mirrored in generate_dita.py.
# -----------------------------------------------------------------------------

def setup_logging(log_path: Path) -> None:
    """Configure dual stdout + per-stage-file logging."""
    root = logging.getLogger()
    for handler in list(root.handlers):
        root.removeHandler(handler)
    formatter = logging.Formatter("%(asctime)s %(levelname)s %(name)s: %(message)s")
    stream = logging.StreamHandler(sys.stdout)
    stream.setLevel(logging.INFO)
    stream.setFormatter(formatter)
    file_handler = logging.FileHandler(log_path, mode="w", encoding="utf-8")
    file_handler.setLevel(logging.DEBUG)
    file_handler.setFormatter(formatter)
    root.setLevel(logging.DEBUG)
    root.addHandler(stream)
    root.addHandler(file_handler)


# -----------------------------------------------------------------------------
# GLC parser (R6, contracts/glc-schema.md)
# -----------------------------------------------------------------------------

@dataclass
class GlcDocument:
    image_filename: str = ""
    bandwidth: str = ""
    bandcentre: str = ""
    warnings: list[str] = field(default_factory=list)


def parse_glc(path: Path) -> GlcDocument:
    """Parse a GLC XML file tolerantly. Never raises.

    Per contracts/glc-schema.md: missing elements yield empty values plus
    a verbatim warning string; malformed XML yields an empty result with
    a single ``"GLC malformed: <reason>"`` warning. Path stripping uses
    ``pathlib.PureWindowsPath(raw).name`` so a Windows ``W:\\foo\\bar.PNG``
    surfaces as ``bar.PNG``.

    The GLC's ``bottom_crop`` is deliberately **not** read: the gram's time
    period (``time_end``) is derived from the pixel height of the referenced
    image, not from the ``.glc`` (issue #148). Many valid image GLCs omit the
    element entirely, so reading it here only produced spurious "invalid GLC"
    warnings.
    """
    doc = GlcDocument()
    try:
        tree = ET.parse(path)
    except ET.ParseError as exc:
        first = str(exc).splitlines()[0] if str(exc) else "unknown error"
        doc.warnings.append(f"GLC malformed: {first}")
        return doc
    except OSError as exc:
        doc.warnings.append(f"GLC malformed: {exc}")
        return doc

    root = tree.getroot()

    filenames = root.findall("data_source/filename")
    if not filenames:
        doc.warnings.append("GLC missing filename")
    else:
        if len(filenames) > 1:
            doc.warnings.append("GLC duplicate filename")
        raw = (filenames[0].text or "").strip()
        if not raw:
            doc.warnings.append("GLC missing filename")
            doc.image_filename = ""
        else:
            doc.image_filename = PureWindowsPath(raw).name

    bandwidth = root.findtext("settings/lofar/bandwidth")
    if bandwidth is None or not bandwidth.strip():
        doc.warnings.append("GLC missing bandwidth")
        doc.bandwidth = ""
    else:
        doc.bandwidth = bandwidth.strip()

    # The frequency band is defined by bandwidth + bandcentre (issue #87):
    # the band spans bandwidth/2 either side of bandcentre. freq_start and
    # freq_end are derived downstream by the generator, not stored here.
    bandcentre = root.findtext("settings/lofar/bandcentre")
    if bandcentre is None or not bandcentre.strip():
        doc.warnings.append("GLC missing bandcentre")
        doc.bandcentre = ""
    else:
        doc.bandcentre = bandcentre.strip()

    return doc


# -----------------------------------------------------------------------------
# Path resolution and walking
# -----------------------------------------------------------------------------

def resolve_glc_path(href: str, content_root: Path, source_dir: Path | None = None) -> Path | None:
    """Resolve a GLC href against the per-gram or per-ten-grams layout (FR-006).

    Returns ``None`` and logs a WARNING when the file cannot be found.
    Hrefs in PowerPoint relationships are URL-encoded (``%20`` for space,
    etc.); the filesystem lookup needs the decoded form.
    """
    if not href:
        LOGGER.warning("GLC not found: empty href")
        return None
    decoded = urllib.parse.unquote(href)
    rel = Path(decoded.replace("\\", "/"))
    candidates: list[Path] = []
    if source_dir is not None:
        candidates.append((source_dir / rel).resolve(strict=False))
    candidates.append((content_root / rel).resolve(strict=False))
    candidates.append((content_root / rel.name).resolve(strict=False))
    for candidate in candidates:
        if candidate.is_file():
            return candidate
    LOGGER.warning("GLC not found: %s", href)
    return None


def _rel_to_root(path: Path, content_root: Path) -> str:
    """Return ``path`` as a POSIX string relative to ``content_root`` when possible."""
    root_abs = content_root.resolve()
    if path.is_relative_to(root_abs):
        return path.relative_to(root_abs).as_posix()
    return path.as_posix()


def resolve_asset_path(href: str, content_root: Path, source_dir: Path | None) -> str:
    """Resolve an asset href (PNG/DOCX/WAV) to a path relative to ``content_root``.

    Returns the resolved POSIX path string when the file is found, or the
    original href stripped of backslashes when it is not — the generator
    treats a missing asset as a dangling reference, not a fatal error.
    """
    if not href:
        return ""
    resolved = resolve_glc_path(href, content_root, source_dir=source_dir)
    if resolved is not None:
        return _rel_to_root(resolved, content_root)
    return href.replace("\\", "/")


def walk_pptxs(input_root: Path, only_subdir: str | None = None) -> Iterator[Path]:
    """Yield every ``.pptx`` under ``input_root`` in deterministic sorted order.

    Skips Office lock files (``~$Foo.pptx``) which Word/PowerPoint
    create alongside any open document and which are not real content.

    When ``only_subdir`` is set, restrict the walk to PPTXs whose path
    (relative to ``input_root``) descends into a folder of that exact
    name as its first segment. This lets the user scope an iteration to
    one chapter for fast debug feedback while keeping ``input_root`` at
    the corpus root, so the resulting CSV's relpaths stay corpus-root-
    relative and dedupe/generate work against the same root unchanged.
    Match is case-insensitive (corpus folder casing varies by author).
    """
    needle = only_subdir.casefold() if only_subdir else None
    for path in sorted(input_root.rglob("*.pptx")):
        if path.name.startswith("~$"):
            continue
        if needle is not None:
            rel_parts = path.relative_to(input_root).parts
            if not rel_parts or rel_parts[0].casefold() != needle:
                continue
        yield path


# -----------------------------------------------------------------------------
# Slug + chapter helpers (R3)
# -----------------------------------------------------------------------------

_SLUG_NON_ALNUM = re.compile(r"[^a-z0-9]+")


def slugify(text: str) -> str:
    """Lower-case, ASCII-only, hyphen-separated slug with collapsed runs."""
    ascii_only = text.encode("ascii", "ignore").decode("ascii").lower()
    slug = _SLUG_NON_ALNUM.sub("-", ascii_only).strip("-")
    return slug


# A ``main`` deck's folder title encodes the teaching week (feature 008).
# "Instructor Week 1 Grams", "Week 01", "Week1" all yield the bare integer
# week number, leading zeros stripped, matched case-insensitively.
_WEEK_TOKEN_RE = re.compile(r"\bweek\s*0*(\d+)\b", re.IGNORECASE)


def week_chapter_number(chapter_title: str | None) -> str:
    """Return the bare-integer week number in ``chapter_title``, or ``""``.

    Used to populate the editable ``target_chapter`` column for ``main`` decks
    so the four week folders (``Week 1`` … ``Week 4``) replace the old
    per-document slicing (feature 008, FR-001). A title with no week token
    (e.g. ``Instructor Pub10_Ed22B_Updated``) returns ``""`` so an analyst can
    assign the week by hand (FR-002).
    """
    if not chapter_title:
        return ""
    match = _WEEK_TOKEN_RE.search(chapter_title)
    return match.group(1) if match else ""


def even_week_assignment(total: int, weeks: int = 4) -> list[str]:
    """Per-gram week labels for ``total`` grams sliced evenly across ``weeks``.

    Feature 009: a no-week ``main`` deck (Pub10, Legacy Pub 10) has no agreed
    per-gram week, so its grams are distributed in **contiguous blocks** —
    ``floor(total/weeks)`` per week with the first ``total % weeks`` weeks each
    taking one extra. Returns a list of length ``total`` of bare-integer week
    strings (``"1"``…``"4"``) in source order, so the first block lands in
    week 1, the next in week 2, and so on. Deterministic and side-effect free.
    """
    base, rem = divmod(total, weeks)
    labels: list[str] = []
    for week in range(1, weeks + 1):
        size = base + (1 if week <= rem else 0)
        labels.extend([str(week)] * size)
    return labels


def deck_target_chapters(
    publication: str, chapter: str | None, gram_count: int,
) -> list[str]:
    """Per-gram ``target_chapter`` values for one deck's grams (feature 009).

    - non-``main`` publication → all empty (no week routing).
    - ``main`` with a ``Week N`` title token → all that week (feature 008).
    - ``main`` with no week token → an **even slice** across the four weeks
      (``even_week_assignment``), replacing the old leave-blank-for-analyst
      path. ``target_chapter`` stays author-editable downstream.
    """
    if publication != "main":
        return [""] * gram_count
    week = week_chapter_number(chapter)
    if week:
        return [week] * gram_count
    return even_week_assignment(gram_count)


def test_number_from_name(name: str) -> int | None:
    """Return the lone integer in a progress-test deck ``name``, or ``None``.

    Progress-test decks are named ``…Progress Test N…`` (the corpus folder /
    file carries a single integer), so ``N`` is the stable publication
    number. Deriving ``progress-test-N`` from the name — rather than the order
    decks happen to be walked — keeps the number identical whether the run
    covers the whole corpus or is scoped to a subset with ``--only``. Leading
    zeros are stripped. Returns ``None`` when the name holds no integer, or
    more than one (ambiguous), leaving the caller to fall back to
    encounter-order allocation.
    """
    numbers = re.findall(r"\d+", name)
    if len(numbers) == 1:
        return int(numbers[0])
    return None


def classify_publication(
    pptx: Path,
    test_pattern: str,
    allocated: dict[str, int],
    final_pattern: str = "",
    final_allocated: dict[str, int] | None = None,
    joining_pattern: str = "",
    joining_allocated: dict[str, int] | None = None,
) -> tuple[str, str | None, str | None]:
    """Return ``(publication, chapter, chapter_slug)`` per R2/R3.

    Progress-test PPTXs are detected by case-insensitive substring match
    against the filename. The test number is taken from the **single integer
    in the deck name** (``test_number_from_name``), so ``progress-test-N`` is
    stable under ``--only`` scoping; a name with no/ambiguous integer falls
    back to stable encounter-order allocation.

    Final-assessment PPTXs (matched against ``final_pattern`` when
    non-empty and a separate ``final_allocated`` map is supplied) get
    their own ``final-assessment-N`` publication prefix. The final
    pattern is checked first so a filename containing both phrases
    routes to the final-assessment bucket.

    Joining-assessment PPTXs (the initial joining assessment, matched
    against ``joining_pattern`` when non-empty with its own
    ``joining_allocated`` map) get a ``joining-assessment-N`` prefix.
    The joining pattern is checked first so a deck deliberately named
    for the joining assessment never falls through to the final or test
    buckets.
    """
    name = pptx.name.lower()
    if joining_pattern and joining_allocated is not None and joining_pattern.lower() in name:
        if pptx.stem not in joining_allocated:
            joining_allocated[pptx.stem] = len(joining_allocated) + 1
        return (f"joining-assessment-{joining_allocated[pptx.stem]}", None, None)
    if final_pattern and final_allocated is not None and final_pattern.lower() in name:
        if pptx.stem not in final_allocated:
            final_allocated[pptx.stem] = len(final_allocated) + 1
        return (f"final-assessment-{final_allocated[pptx.stem]}", None, None)
    if test_pattern.lower() in name:
        number = test_number_from_name(pptx.stem)
        if number is None:
            if pptx.stem not in allocated:
                allocated[pptx.stem] = len(allocated) + 1
            number = allocated[pptx.stem]
        return (f"progress-test-{number}", None, None)
    chapter_title = pptx.parent.name
    return ("main", chapter_title, slugify(chapter_title))


# -----------------------------------------------------------------------------
# Shape grouping stub (FR-015 / R1)
# -----------------------------------------------------------------------------

@dataclass
class GlcLink:
    display_text: str
    href: str


@dataclass
class GramPlaceholder:
    gram_id: str
    vessel_name: str
    png_href: str | None
    glc_links: list[GlcLink]


def _label_ends_in_questions(label: str) -> bool:
    """True when ``label``'s final word is "questions" (case-insensitive).

    Matches "N Questions", "Questions", "10 questions" and tolerates
    trailing whitespace/punctuation; does not match "acquisitions".
    """
    return re.search(
        rf"\b{QUESTIONS_LABEL_WORD}\b\W*$", label, re.IGNORECASE
    ) is not None


def _shape_level_hyperlink(shape) -> str | None:
    """Return the shape-level hyperlink target, or None.

    Walks the lxml element directly so we don't depend on python-pptx
    exposing a high-level accessor (it doesn't, for shape-level clicks).
    Searches every descendant ``p:cNvPr`` so the lookup works regardless
    of the shape's outer XML wrapper — ``p:sp`` (autoshape/textbox),
    ``p:pic`` (picture), ``p:cxnSp`` (connector), or
    ``p:graphicFrame`` (chart/table/etc.). The actual ``a:hlinkClick``
    always sits under ``cNvPr`` regardless of wrapper.
    """
    for c_nv_pr in shape._element.iter(qn("p:cNvPr")):
        hlink = c_nv_pr.find(qn("a:hlinkClick"))
        if hlink is None:
            continue
        rel_id = hlink.get(qn("r:id"))
        if not rel_id:
            continue
        try:
            return shape.part.rels[rel_id].target_ref
        except KeyError:
            continue
    return None


def _run_hyperlinks_in_shape(shape) -> list[tuple[str, str]]:
    """Return a list of ``(visible_text, href)`` per hyperlinked text run.

    When a paragraph contains exactly one hyperlinked run the whole
    paragraph's combined text is returned as ``visible_text`` — legacy
    decks frequently split a single label across multiple runs
    (e.g. ``"Lofar"`` carries the hyperlink, a sibling run carries
    ``" 2"``), and the user-visible label is the paragraph total.

    When a paragraph carries multiple hyperlinks (the multi-channel
    "Lofar box" pattern), each link keeps its own run-text as label.
    """
    pairs: list[tuple[str, str]] = []
    if not shape.has_text_frame:
        return pairs
    for paragraph in shape.text_frame.paragraphs:
        link_runs: list[tuple[object, str]] = []
        for run in paragraph.runs:
            try:
                href = run.hyperlink.address
            except Exception:
                href = None
            if href:
                link_runs.append((run, href))
        if not link_runs:
            continue
        if len(link_runs) == 1:
            para_text = "".join(r.text or "" for r in paragraph.runs)
            _, href = link_runs[0]
            pairs.append((para_text, href))
        else:
            for run, href in link_runs:
                pairs.append((run.text or "", href))
    return pairs


def _gram_folder_key(href: str) -> str:
    """Last directory segment of a hyperlink target, URL-decoded and
    lowercased. Used to associate each .glc with the gram-header that
    points at the same ``GramNN/`` directory on disk.

    Returns the empty string if the href has no parent directory
    (e.g. a bare filename) — such links can't be folder-matched and
    are skipped by the caller.
    """
    if not href:
        return ""
    decoded = urllib.parse.unquote(href.split("?", 1)[0].split("#", 1)[0])
    path = PurePosixPath(decoded.replace("\\", "/"))
    parts = path.parts
    if len(parts) < 2:
        return ""
    return parts[-2].lower()


def _bbox(shape) -> tuple[int, int, int, int]:
    """Return ``(left, top, right, bottom)`` in EMUs."""
    left = shape.left or 0
    top = shape.top or 0
    width = shape.width or 0
    height = shape.height or 0
    return (left, top, left + width, top + height)


def _iter_leaf_shapes(shapes):
    """Yield every leaf shape on a slide, expanding GROUPs (shape_type=6)."""
    for shape in shapes:
        if getattr(shape, "shape_type", None) == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_leaf_shapes(shape.shapes)
        else:
            yield shape


def _slide_diagram_hyperlinks(slide) -> list[tuple[str, str]]:
    """Return ``(display_text, href)`` for every external hyperlink
    reachable from a SmartArt diagram on ``slide``.

    SmartArt nodes can carry click hyperlinks stored in the diagram's
    own relationships (``ppt/diagrams/_rels/data1.xml.rels``,
    ``drawing1.xml.rels``, …), invisible to python-pptx's slide-level
    shape walks. Walks from the slide's part through any diagram
    relationships and recurses one level to cover both data1 and
    drawing1 hyperlink sets. Targets are deduplicated by URI.

    ``display_text`` is harvested from the SmartArt data tree: each
    ``<dgm:pt>`` carries its own ``<dgm:prSet><a:hlinkClick r:id="…"/>``
    plus inline ``<a:t>`` runs. We index node text by rId, then look it
    up when the matching relationship is found.
    """
    DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
    A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
    R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    PT_TAG = f"{{{DGM_NS}}}pt"
    HLINK_TAG = f"{{{A_NS}}}hlinkClick"
    T_TAG = f"{{{A_NS}}}t"
    RID_ATTR = f"{{{R_NS}}}id"

    pairs: list[tuple[str, str]] = []
    seen_targets: set[str] = set()
    visited: set[int] = set()

    def _is_diagram_part(part) -> bool:
        try:
            return "/diagrams/" in str(part.partname).lower()
        except Exception:
            return False

    def _text_by_rid(part) -> dict[str, str]:
        """Map rId → concatenated visible text from each SmartArt node."""
        result: dict[str, str] = {}
        try:
            tree = ET.fromstring(part.blob)
        except Exception:
            return result
        for pt in tree.iter(PT_TAG):
            rid = None
            for hlink in pt.iter(HLINK_TAG):
                rid = hlink.get(RID_ATTR)
                if rid:
                    break
            if not rid:
                continue
            text = "".join((t.text or "") for t in pt.iter(T_TAG)).strip()
            result[rid] = text
        return result

    def _walk(part) -> None:
        if id(part) in visited:
            return
        visited.add(id(part))
        node_text = _text_by_rid(part)
        for rid, rel in part.rels.items():
            if rel.is_external:
                if rel.reltype.endswith("/hyperlink"):
                    target = rel.target_ref
                    if target and target not in seen_targets:
                        seen_targets.add(target)
                        pairs.append((node_text.get(rid, ""), target))
                continue
            try:
                target_part = rel.target_part
            except Exception:
                continue
            if _is_diagram_part(target_part):
                _walk(target_part)

    for rel in slide.part.rels.values():
        if rel.is_external:
            continue
        try:
            target_part = rel.target_part
        except Exception:
            continue
        if _is_diagram_part(target_part):
            _walk(target_part)

    return pairs


def is_framing_slide(slide) -> bool:
    """Return True for welcome / exit slides that should be skipped.

    Detection is by title text prefix (``"Welcome to "`` / ``"End of "``)
    so it survives slide-position shuffles and is robust against real
    decks that may add or remove framing slides.
    """
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = (shape.text_frame.text or "").strip()
        if not text:
            continue
        if text.startswith(FRAMING_TITLE_PREFIXES):
            return True
    return False


def _descriptor_text(header) -> str:
    """Collapse a header shape's text runs into its ``"Gram N: <detail>"`` caption."""
    return "".join(
        run.text or "" for para in header.text_frame.paragraphs for run in para.runs
    ).strip()


def _reconnect_orphan_glc(
    glc_href: str,
    folder_key: str,
    num_to_header: dict,
    content_root: Path | None,
    source_dir: Path | None,
):
    """Reconnect a folder-orphaned ``.glc`` to its gram header via a sibling
    analysis sheet, or return ``None`` when it can't be confidently recovered.

    A ``.glc`` whose ``GramNN/`` folder matches no header on the slide is
    usually a gram whose *analysis-sheet* header hyperlink was left pointing at
    an earlier build's folder (a stale link), so the folder-key pairing in
    ``extract_grams_from_slide`` can't associate them. When the ``.glc``'s own
    on-disk folder holds an analysis-sheet image — proof this is a real gram,
    and the very sheet ``generate_dita`` will surface once the row exists (see
    the Lofar-folder recovery in ``gram_to_rows``) — reconnect it to the header
    carrying the same gram number rather than dropping it and warning.

    Returns that header, or ``None`` when recovery isn't safe: no filesystem
    context, a folder key with no gram number, the ``.glc`` doesn't resolve on
    disk, no sibling analysis sheet beside it, or no number-matched header. In
    every ``None`` case the caller warns and drops exactly as before.
    """
    if content_root is None:
        return None
    if not re.search(r"\d", folder_key):
        return None
    resolved = resolve_glc_path(glc_href, content_root, source_dir=source_dir)
    if resolved is None:
        return None
    if _recover_analysis_sheet([resolved.parent], content_root) is None:
        return None
    return num_to_header.get(_gram_num_from_id(folder_key))


def extract_grams_from_slide(
    slide,
    slide_num: int,
    content_root: Path | None = None,
    source_dir: Path | None = None,
) -> list[GramPlaceholder]:
    """Return the gram placeholders on ``slide`` per reverse-spec §4.

    Each gram tile has:

    - a header shape carrying a *shape-level* hyperlink to an Analysis
      Sheet (``.docx`` or ``.png``); its text is a ``"Gram N: <detail>"``
      descriptor (split at the first colon — left side is the
      student-visible label, right side is instructor-visible);
    - one or more text-frame shapes positioned beneath the header
      containing text *runs* hyperlinked to ``.glc`` files in the same
      gram folder. Every Lofar text-run hyperlink in the audited corpus
      targets a ``.glc``; the ``.wav`` case is one indirection deeper,
      inside the ``.glc``'s ``data_source/filename`` element (see
      ``high-level-spec.md`` §1.5).

    Grouping is done by spatial proximity: for each header we collect
    every Lofar text run carried by text-frame shapes whose bounding
    box sits below the header and overlaps its horizontal extent. This
    survives renamed shapes, ad-hoc shape ordering, and GROUP wrappers.
    """
    leaves = list(_iter_leaf_shapes(slide.shapes))

    # 1) Identify headers (shape-level hyperlink to an analysis-sheet
    #    asset, on a text-bearing shape).
    headers: list[tuple[object, str]] = []  # (shape, analysis_href)
    for shape in leaves:
        href = _shape_level_hyperlink(shape)
        if href is None:
            continue
        if not shape.has_text_frame:
            continue
        # Skip vestigial overlay shapes whose shape-level hyperlinks point
        # at absolute file:/// URIs from a long-gone authoring path. The
        # live header buttons use paths relative to the PPTX; the dead
        # overlay (e.g. Group 197/Rectangle children) is left over from
        # earlier authoring iterations and never resolves.
        if href.lower().startswith("file:///"):
            continue
        # Whitelist the analysis-sheet extensions. Legacy authoring
        # sometimes promotes a .glc text-run hyperlink to the shape
        # level on the small .glc-bearing shapes; without this filter
        # those would be mistaken for gram headers.
        if not href.lower().endswith(ANALYSIS_SHEET_EXTENSIONS):
            continue
        # Skip self-test navigation buttons labelled "N Questions" — they
        # link to an image but are not grams (see QUESTIONS_LABEL_WORD).
        label = " ".join((shape.text_frame.text or "").split())
        if _label_ends_in_questions(label):
            LOGGER.info(
                "Slide %d: ignoring %r link (label ends in %r) — not a gram",
                slide_num, label, QUESTIONS_LABEL_WORD,
            )
            continue
        headers.append((shape, href))

    # 2) Identify candidate Lofar shapes — shapes whose hyperlink targets
    #    a .glc. Two authoring styles both occur in the legacy corpus:
    #      (a) text-bearing shapes (autoshape/textbox) with the .glc
    #          hyperlink on a text *run* — the common case;
    #      (b) picture shapes with the .glc hyperlink at the shape level
    #          — used for at least one gram per slide in some decks
    #          (e.g. the gram whose analysis sheet is "V III .doc").
    candidates: list[tuple[object, list[tuple[str, str]]]] = []
    header_ids = {id(s) for s, _ in headers}
    for shape in leaves:
        if id(shape) in header_ids:
            continue
        # Fallback display text: legacy decks frequently put the hyperlink
        # on a zero-width run while the visible label sits in a sibling
        # run of the same shape. When the run.text is empty, use the
        # whole shape's collapsed text-frame text as the link label.
        shape_text = ""
        if getattr(shape, "has_text_frame", False):
            shape_text = " ".join((shape.text_frame.text or "").split())
        keep: list[tuple[str, str]] = []
        # (a) text-run .glc hyperlinks
        for text, href in _run_hyperlinks_in_shape(shape):
            if href.lower().endswith(".glc"):
                display = " ".join((text or "").split()) or shape_text
                keep.append((display, href))
            else:
                LOGGER.warning(
                    "Lofar text run hyperlinks to a non-.glc target (%r); "
                    "expected .glc — row dropped", href,
                )
        # (b) shape-level .glc hyperlink (picture-style). Skip vestigial
        # absolute file:/// URIs the same way the header step does.
        shape_href = _shape_level_hyperlink(shape)
        if (
            shape_href
            and shape_href.lower().endswith(".glc")
            and not shape_href.lower().startswith("file:///")
        ):
            keep.append((shape_text, shape_href))
        if keep:
            candidates.append((shape, keep))

    # 2b) SmartArt-embedded .glc hyperlinks. Pulled from the diagram's
    #     own .rels files (``ppt/diagrams/_rels/...``) rather than the
    #     slide's, so they're invisible to the per-shape walk above.
    #     At least one gram per slide in the audited corpus is
    #     authored as SmartArt (e.g. the gram whose analysis sheet is
    #     "V III .doc"). Attached to the candidates list as a synthetic
    #     entry (no associated shape) — the folder-key pairing in
    #     step 3 doesn't reference the shape.
    diagram_keep: list[tuple[str, str]] = []
    for text, href in _slide_diagram_hyperlinks(slide):
        if not href.lower().endswith(".glc"):
            continue
        if href.lower().startswith("file:///"):
            continue
        diagram_keep.append((text, href))
    if diagram_keep:
        candidates.append((None, diagram_keep))

    # 3) Match each .glc to its gram header by shared parent folder on
    #    disk (e.g. both targets live under ``.../Gram001/``). This is
    #    far more robust than spatial proximity: it survives shapes
    #    moved off-screen, off-grid layouts, hidden overlay rectangles,
    #    and decks where the .glc shapes don't sit directly beneath
    #    their header rectangle. Per the corpus survey, every live
    #    header and its .glc children share a single ``GramNN/`` parent.
    folder_to_header: dict[str, tuple[object, str]] = {}
    for header, href in headers:
        key = _gram_folder_key(href)
        if not key:
            h_left, h_top, _, _ = _bbox(header)
            LOGGER.warning(
                "Slide %d: header at top=%d left=%d has no gram-folder in href %r",
                slide_num, h_top, h_left, href,
            )
            continue
        if key in folder_to_header:
            LOGGER.warning(
                "Slide %d: duplicate gram folder %r — second header for %r ignored",
                slide_num, key, href,
            )
            continue
        folder_to_header[key] = (header, href)

    # Index headers by their parsed gram number so a .glc whose gram folder
    # matches no header — because that gram's analysis-sheet header links to a
    # stale earlier-build folder — can still be reconnected to its gram via a
    # sibling analysis sheet (see _reconnect_orphan_glc). First header wins on a
    # numeric collision; numberless headers can't disambiguate, so are skipped.
    num_to_header: dict[str, object] = {}
    for header, _analysis_href in headers:
        h_gram_id, _ = _split_descriptor(_descriptor_text(header))
        if re.search(r"\d", h_gram_id):
            num_to_header.setdefault(_gram_num_from_id(h_gram_id), header)

    header_to_pairs: dict[int, list[tuple[str, str]]] = defaultdict(list)
    for cand, pairs in candidates:
        for text, glc_href in pairs:
            key = _gram_folder_key(glc_href)
            if not key:
                LOGGER.warning(
                    "Slide %d: .glc %r has no gram-folder in href",
                    slide_num, glc_href,
                )
                continue
            hdr_entry = folder_to_header.get(key)
            if hdr_entry is None and len(key) > 1 and key.endswith("a"):
                # Legacy authoring pattern: gram's .glc files live in a
                # sibling folder named with a trailing 'a' suffix
                # (e.g. analysis sheet in "Gram_11/" but .glc in
                # "Gram_11a/"). Try the trimmed key as a fallback so
                # the gram still groups correctly.
                fallback = key[:-1]
                hdr_entry = folder_to_header.get(fallback)
                if hdr_entry is not None:
                    LOGGER.info(
                        "Slide %d: .glc %r folder %r matched header %r via "
                        "trailing-'a' fallback", slide_num, glc_href, key, fallback,
                    )
            if hdr_entry is None:
                recovered_header = _reconnect_orphan_glc(
                    glc_href, key, num_to_header, content_root, source_dir,
                )
                if recovered_header is not None:
                    # A stale analysis-sheet link left this .glc folder-orphaned,
                    # but its sibling analysis sheet proves the gram is real —
                    # reconnect it by gram number rather than warn-and-drop.
                    LOGGER.info(
                        "Slide %d: .glc %r names folder %r with no folder-matched "
                        "header (its analysis link points elsewhere); reconnected "
                        "to the 'Gram %s' header via its sibling analysis sheet",
                        slide_num, glc_href, key, _gram_num_from_id(key),
                    )
                    header_to_pairs[id(recovered_header)].append((text, glc_href))
                    continue
                LOGGER.warning(
                    "Slide %d: .glc %r names folder %r but no matching header on this slide",
                    slide_num, glc_href, key,
                )
                continue
            header_id = id(hdr_entry[0])
            header_to_pairs[header_id].append((text, glc_href))

    # 3b) Optional filesystem validation. Legacy decks can carry
    #     .glc hyperlinks whose target no longer exists on disk —
    #     usually a renamed/removed channel that the author forgot to
    #     unlink. Such links don't click anywhere in PowerPoint either,
    #     so we drop them rather than emit a phantom Lofar row. The
    #     check only runs when the caller provides ``content_root``;
    #     unit tests that don't have a filesystem layout in mind keep
    #     working unchanged.
    if content_root is not None:
        dropped = 0
        for header_id in list(header_to_pairs.keys()):
            kept: list[tuple[str, str]] = []
            for text, href in header_to_pairs[header_id]:
                resolved = resolve_glc_path(href, content_root, source_dir=source_dir)
                if resolved is None:
                    LOGGER.warning(
                        "Slide %d: .glc target %r not present on disk — "
                        "dropped as legacy artefact", slide_num, href,
                    )
                    dropped += 1
                    continue
                kept.append((text, href))
            header_to_pairs[header_id] = kept
        if dropped:
            LOGGER.info(
                "Slide %d: filesystem-validation dropped %d stale .glc link(s)",
                slide_num, dropped,
            )

    # 3c) Per-header dedup. Legacy decks sometimes carry a second
    #     hyperlink to the same .glc within a gram — typically a
    #     leftover from iterative authoring, recognisable by an
    #     integer-only label (e.g. "1", "2") next to the real label
    #     ("0-300 Hz"). Keep the entry with the longest display text
    #     per unique href; first-seen order is preserved.
    deduped = 0
    for header_id in list(header_to_pairs.keys()):
        best_by_href: dict[str, tuple[str, str]] = {}
        order: list[str] = []
        for text, href in header_to_pairs[header_id]:
            existing = best_by_href.get(href)
            if existing is None:
                best_by_href[href] = (text, href)
                order.append(href)
            elif len(text) > len(existing[0]):
                best_by_href[href] = (text, href)
                deduped += 1
            else:
                deduped += 1
        header_to_pairs[header_id] = [best_by_href[h] for h in order]
    if deduped:
        LOGGER.info(
            "Slide %d: deduplicated %d .glc link(s) — kept longest-label variant",
            slide_num, deduped,
        )

    # 4) Build grams in reading order (top, left) for stable CSV output.
    grams: list[GramPlaceholder] = []
    headers.sort(key=lambda hb: (hb[0].top or 0, hb[0].left or 0))

    for header, analysis_href in headers:
        lofar_pairs = header_to_pairs.get(id(header), [])

        gram_id, instructor_detail = _split_descriptor(_descriptor_text(header))

        # Deleted-gram remnant: when a gram is removed in PowerPoint the
        # emptied header button often survives, still carrying its
        # shape-level hyperlink to the old analysis sheet but with its
        # "Gram N: ..." caption gone. That shape passes the header test
        # above yet yields an empty gram_id, and its now-orphaned GramNN/
        # folder no longer has matching .glc shapes, so it would emit a
        # lone analysis row with a blank gram_id -- an invalid row that
        # dedupe rejects (require_field). A header with neither a gram
        # number nor any Lofar link is not publishable content, so drop
        # it here (loudly) rather than the operator hand-deleting the row
        # on every run. A header that *does* carry Lofar links but no
        # parseable number is a different problem (a caption we failed to
        # read) -- keep it and let the no-Lofar/number warnings surface.
        if not gram_id and not lofar_pairs:
            h_left, h_top, _, _ = _bbox(header)
            LOGGER.warning(
                "Slide %d: skipping deleted-gram remnant at top=%d left=%d "
                "(analysis href %r) -- no gram number and no Lofar links",
                slide_num, h_top, h_left, analysis_href,
            )
            continue

        if not lofar_pairs:
            h_left, h_top, _, _ = _bbox(header)
            LOGGER.warning(
                "Slide %d: gram header at top=%d left=%d has no Lofar box",
                slide_num, h_top, h_left,
            )

        glc_links = [GlcLink(display_text=t.strip(), href=h) for t, h in lofar_pairs]
        grams.append(GramPlaceholder(
            gram_id=gram_id,
            vessel_name=instructor_detail,
            png_href=analysis_href,
            glc_links=glc_links,
        ))

    # Sort by leading-integer gram_id so "2" precedes "10" (lexicographic
    # would reverse them). Falls back to vessel name when ids are
    # missing or equal. Determines both CSV row order and the report's
    # per-gram listing — single source of truth for ordering.
    def _gram_sort_key(g: GramPlaceholder) -> tuple[float, str]:
        m = re.match(r"\d+", g.gram_id or "")
        num = float(m.group(0)) if m else float("inf")
        return (num, (g.vessel_name or "").lower())
    grams.sort(key=_gram_sort_key)
    return grams


def _split_descriptor(descriptor: str) -> tuple[str, str]:
    """Split a `"Gram N: <detail>"` descriptor at the first colon.

    Returns ``(gram_id, instructor_detail)``. If no colon is present, the
    whole text is treated as the gram_id and instructor_detail is empty.
    Gram_id is normalised to a plain integer string (``"7"``, ``"12"``)
    when the left side parses as ``"Gram N"`` — the canonical CSV form
    per csv-schema.md so authors can renumber with a bare number when
    refactoring content between chapters.
    """
    # Collapse runs of whitespace (legacy decks pad gram titles with
    # multi-space sequences to force in-shape line breaks).
    descriptor = " ".join(descriptor.split())
    if not descriptor:
        return ("", "")
    left, sep, right = descriptor.partition(":")
    left = left.strip()
    right = right.strip() if sep else ""
    m = re.match(r"^Gram\s+(\d+)$", left, re.IGNORECASE)
    if m:
        left = str(int(m.group(1)))
    return (left, right)


# -----------------------------------------------------------------------------
# Row construction
# -----------------------------------------------------------------------------

def _gram_num_from_id(gram_id: str) -> str:
    digits = re.findall(r"\d+", gram_id)
    return digits[0] if digits else "00"


def _png_file_size(png_path: str, content_root: Path) -> str:
    """Return the on-disk size of ``png_path`` (as a decimal string) or ``""``.

    ``png_path`` is recorded as a content-root-relative POSIX string;
    resolve it back to a real Path and stat it. Used by the duplicate-
    detection workflow (the author groups duplicate lofars/analysis assets
    by exact byte size rather than filename, since names can drift).
    """
    if not png_path:
        return ""
    candidate = (content_root / png_path)
    try:
        return str(candidate.stat().st_size)
    except OSError:
        return ""


def _jpeg_pixel_height(fh) -> int | None:
    """Walk a JPEG's segment markers and return its pixel height, or ``None``.

    ``fh`` is an open binary file positioned anywhere; the scan restarts at the
    first marker after the SOI. Reads only segment headers (never the entropy-
    coded scan data), skipping non-frame segments by their declared length until
    a Start-Of-Frame marker is reached, whose body carries ``precision(1),
    height(2), width(2)``. Returns ``None`` on a truncated or unexpected stream.
    """
    fh.seek(2)  # past the SOI (0xFFD8)
    # Frame headers that carry image dimensions. DHT (0xC4), JPG (0xC8) and
    # DAC (0xCC) share the 0xCn range but are not frame headers, so they are
    # skipped like any other length-prefixed segment.
    sof_markers = frozenset((
        0xC0, 0xC1, 0xC2, 0xC3, 0xC5, 0xC6, 0xC7,
        0xC9, 0xCA, 0xCB, 0xCD, 0xCE, 0xCF,
    ))
    while True:
        prefix = fh.read(1)
        if not prefix:
            return None
        if prefix != b"\xff":
            continue
        marker = fh.read(1)
        while marker == b"\xff":  # fill bytes before the marker code
            marker = fh.read(1)
        if not marker:
            return None
        code = marker[0]
        # Standalone markers (no length payload): SOI, EOI, RSTn, TEM.
        if code == 0xD8 or code == 0xD9 or 0xD0 <= code <= 0xD7 or code == 0x01:
            continue
        length_bytes = fh.read(2)
        if len(length_bytes) < 2:
            return None
        seg_len = int.from_bytes(length_bytes, "big")
        if seg_len < 2:
            return None
        if code in sof_markers:
            body = fh.read(seg_len - 2)
            if len(body) < 5:
                return None
            return int.from_bytes(body[1:3], "big")
        fh.seek(seg_len - 2, 1)  # skip this segment's body


def _image_pixel_height(png_path: str, content_root: Path) -> int | None:
    """Return the pixel height (scan-line count) of an image, or ``None``.

    Issue #148: a pre-rendered gram's time period is derived from the image
    itself — the number of horizontal scan lines, i.e. the pixel height — not
    from any value in the ``.glc`` (the legacy viewer multiplies the scan-line
    count by an update period that is always ``1`` second, so seconds == pixel
    rows). Reads only the header of PNG, JPEG and GIF files using the standard
    library, keeping Pillow off the runtime path per the air-gap dependency
    budget. Returns ``None`` when the file is absent, truncated, or not a
    recognised/parsable format, so the caller leaves ``time_end`` blank and the
    usual missing-view-field handling (warn, or ``--relaxed`` default) applies.
    """
    if not png_path:
        return None
    candidate = content_root / png_path
    try:
        with candidate.open("rb") as fh:
            head = fh.read(26)
            if head[:8] == b"\x89PNG\r\n\x1a\n":
                # PNG mandates IHDR as the first chunk; height is the four
                # big-endian bytes at offset 20 (8 sig + 4 len + 4 tag + 4 w).
                return int.from_bytes(head[20:24], "big") if len(head) >= 24 else None
            if head[:6] in (b"GIF87a", b"GIF89a"):
                # Logical-screen height: two little-endian bytes at offset 8.
                return int.from_bytes(head[8:10], "little") if len(head) >= 10 else None
            if head[:2] == b"\xff\xd8":
                return _jpeg_pixel_height(fh)
    except OSError:
        return None
    return None


def _asset_file_missing(png_path: str, content_root: Path) -> bool:
    """True when a target asset is *expected* (non-empty path) but absent on disk.

    A blank ``png_path`` means the row carries no asset at all (an assetless
    GLC link, or a gram with no analysis hyperlink) — that is not "missing",
    just legitimately empty, so it returns ``False``. Only a populated
    reference whose file cannot be found on disk counts, so the caller can
    flag it for the author to fix or drop (mirrors the resolution
    ``_png_file_size`` performs, but reports presence rather than size).
    """
    if not png_path:
        return False
    return not (content_root / png_path).is_file()


def _gram_folder_missing(
    gram: GramPlaceholder, content_root: Path, source_dir: Path | None,
) -> bool:
    """True when *every* asset folder this gram references is absent on disk.

    A gram whose whole per-gram folder an author has deleted references its
    assets — the analysis-sheet image and the Lofar ``.glc`` files — out of
    that now-missing directory. This is deliberately distinct from a single
    deleted *file* with the folder still present: the latter stays flagged
    (``ASSET_MISSING_WARNING``) so the author restores the one file, whereas a
    vanished folder means the gram itself is gone, so the caller drops it from
    the CSV entirely rather than emit a topic full of dangling references (that
    would only fail later in Oxygen/DITA-OT).

    Fires only when the gram carries at least one href with a directory
    component *and none* of those directories resolves on disk (checked under
    the deck folder then the content root, mirroring ``resolve_glc_path``'s
    candidate order). A gram whose hrefs are all bare filenames — no folder to
    test — or any one of whose asset folders still exists returns ``False``, so
    a merely-deleted individual file continues to be flagged, not dropped.
    """
    hrefs = [link.href for link in gram.glc_links if link.href]
    if gram.png_href:
        hrefs.append(gram.png_href)
    referenced_dirs: list[Path] = []
    for href in hrefs:
        decoded = urllib.parse.unquote(href.split("?", 1)[0].split("#", 1)[0])
        rel = Path(decoded.replace("\\", "/"))
        if rel.parent == Path("."):
            continue  # bare filename — no folder to test
        referenced_dirs.append(rel.parent)
    if not referenced_dirs:
        return False
    for rel_dir in referenced_dirs:
        candidates: list[Path] = []
        if source_dir is not None:
            candidates.append(source_dir / rel_dir)
        candidates.append(content_root / rel_dir)
        if any(candidate.is_dir() for candidate in candidates):
            return False
    return True


def _is_analysis_sheet_image(name: str) -> bool:
    """True when ``name`` looks like an analysis-sheet image file.

    Matches the ``analysis`` token (or a known misspelling, e.g. ``analaysis``)
    as a case-insensitive substring of the stem, with an image extension —
    mirroring ``snapshot_analysis_docs._is_analysis_name`` so the Lofar-folder
    fallback recognises the same sheets the snapshotter renders.
    """
    p = PurePosixPath(name)
    if p.suffix.lower() not in ANALYSIS_IMAGE_EXTENSIONS:
        return False
    stem = p.stem.lower()
    if ANALYSIS_NAME_TOKEN in stem:
        return True
    return any(typo in stem for typo in ANALYSIS_NAME_MISSPELLINGS)


def _recover_analysis_sheet(
    lofar_dirs: list[Path], content_root: Path,
) -> str | None:
    """Find an analysis-sheet image sitting in a gram's Lofar folder(s).

    When a gram's analysis hyperlink is stale — a legacy href left pointing at
    an earlier build's folder — the real analysis sheet often still sits beside
    the gram's Lofar ``.glc`` files. Search each resolved Lofar directory (in
    the order the Lofars appear, so the result is deterministic) for an analysis
    image and return the first match as a ``content_root``-relative POSIX path,
    or ``None`` when none is found. Directories are scanned non-recursively and
    their entries sorted so a folder holding two candidate images always yields
    the same one (the determinism/idempotency invariant).
    """
    for lofar_dir in lofar_dirs:
        if not lofar_dir.is_dir():
            continue
        for entry in sorted(lofar_dir.iterdir()):
            if entry.is_file() and _is_analysis_sheet_image(entry.name):
                return _rel_to_root(entry, content_root)
    return None


# A demon marker's stem is exactly ``demon`` (the first) or ``demon-N`` (the
# Nth, N>=2). Anchored so ``demonstrate.glc`` and other ``demon``-prefixed files
# never match. Case-insensitive.
_DEMON_MARKER_RE = re.compile(r"^demon(?:-(\d+))?$", re.IGNORECASE)


def _find_demon_markers(lofar_dir: Path) -> list[Path]:
    """Return the ``demon.glc`` / ``demon-N.glc`` markers in ``lofar_dir``.

    Ordered ``demon`` (1) then ``demon-2``, ``demon-3``, … by their numeric
    suffix — a plain lexical sort would wrongly put ``demon-2.glc`` before
    ``demon.glc`` (``-`` < ``.``). Deterministic so multiple demons emit in a
    stable order ahead of the Lofars.
    """
    markers: list[tuple[int, Path]] = []
    for entry in sorted(lofar_dir.iterdir()):
        if not (entry.is_file() and entry.suffix.lower() == ".glc"):
            continue
        match = _DEMON_MARKER_RE.match(entry.stem)
        if match is None:
            continue
        order = int(match.group(1)) if match.group(1) else 1
        markers.append((order, entry))
    return [entry for _order, entry in sorted(markers, key=lambda t: (t[0], t[1].name))]


def demon_rows_for_gram(
    gram: GramPlaceholder,
    lofar_dirs: list[Path],
    publication: str,
    chapter: str | None,
    content_root: Path,
    gram_num: str,
    target_doc: str,
    target_chapter: str,
    relaxed: bool,
) -> list[dict]:
    """Emit one ``topic_type="demon"`` row per demon marker for this gram (#151).

    A demon image is an alternately-rendered gram view that leads the gram's
    page. It is signalled by a ``demon.glc`` marker ingest writes into the gram
    folder (never hyperlinked from the slide), so it is discovered here by
    scanning the folder of the gram's **first resolved Lofar ``.glc``**
    (``lofar_dirs[0]``). A gram with no resolved Lofar has no folder to inspect
    and gets no demon row. Each marker becomes one demon row ordered ahead of
    the Lofar rows; the demon's ``time_end`` is the image's pixel height (issue
    #148) and its band is the fixed 0 - 40 Hz read from the marker.
    """
    if not lofar_dirs:
        return []
    lofar_dir = lofar_dirs[0]
    if not lofar_dir.is_dir():
        return []

    rows: list[dict] = []
    for i, marker in enumerate(_find_demon_markers(lofar_dir), start=1):
        warnings: list[str] = []
        glc = parse_glc(marker)
        warnings.extend(glc.warnings)
        glc_path = _rel_to_root(marker, content_root)
        png_path = ""
        time_end = ""
        if glc.image_filename:
            png_path = resolve_asset_path(
                glc.image_filename, content_root, source_dir=marker.parent)
            if Path(png_path).suffix.lower() in GRAMFRAME_GLC_EXTENSIONS:
                height = _image_pixel_height(png_path, content_root)
                if height is not None:
                    time_end = str(height)

        row = {
            "publication": publication,
            "chapter": chapter or "",
            "target_doc": target_doc,
            "target_chapter": target_chapter,
            "gram_id": gram.gram_id,
            "vessel_name": gram.vessel_name,
            "topic_type": "demon",
            "sequence": str(i),
            "topic_filename": f"gram_{gram_num}.dita",
            "display_text": "",
            "link_href": "",
            "glc_path": glc_path,
            "time_end": time_end,
            "bandwidth": glc.bandwidth,
            "bandcentre": glc.bandcentre,
            "png_path": png_path,
            "target_ext": Path(png_path).suffix.lower(),
            "file_size": _png_file_size(png_path, content_root),
            "wav_treatment": "",
            "warnings": "",
        }
        asset_missing = _asset_file_missing(png_path, content_root)
        # A demon is always rendered as an inline image GramFrame, so it needs
        # the same view fields an image Lofar does. ``time_end`` is image-derived
        # (dangles when the image is absent, like every image row); the band
        # comes from the demon marker ingest wrote (0 - 40 Hz). Mirror the image
        # Lofar handling: default under --relaxed, else warn (never fail-fast --
        # a demon row is topic_type="demon", exempt from glc_view_problems).
        if row["target_ext"] in GRAMFRAME_GLC_EXTENSIONS:
            for field_name in GLC_VIEW_FIELDS:
                if (row[field_name] or "").strip():
                    continue
                if relaxed:
                    row[field_name] = RELAXED_DEFAULT
                    warnings.append(
                        f"demon missing {field_name} — defaulted to "
                        f"{RELAXED_DEFAULT} (--relaxed)")
                elif field_name == "time_end":
                    if not asset_missing:
                        warnings.append(
                            "demon time period unknown — could not read image "
                            "height")
                else:
                    warnings.append(
                        f"demon missing {field_name} — GramFrame cannot render")
        if asset_missing:
            warnings.append(ASSET_MISSING_WARNING)
        row["warnings"] = ", ".join(warnings)
        rows.append(row)
    return rows


def gram_to_rows(
    gram: GramPlaceholder,
    publication: str,
    chapter: str | None,
    chapter_slug: str | None,
    content_root: Path,
    source_dir: Path,
    target_doc: str = "",
    target_chapter: str = "",
    relaxed: bool = False,
) -> list[dict]:
    """Expand one gram into N+1 CSV rows (N GLC links + 1 analysis row).

    ``target_chapter`` is the editable routing target — for ``main`` decks this
    is the bare-integer week number (feature 008); empty falls back to the
    immutable source ``chapter`` downstream.

    ``relaxed`` substitutes ``RELAXED_DEFAULT`` for any blank GLC view field
    (issue #92) so the toolchain can run against an incomplete corpus; the strict
    default leaves the blank in place for ``main`` to hard-fail on.
    """
    rows: list[dict] = []
    gram_num = _gram_num_from_id(gram.gram_id)
    # Folders holding this gram's resolved Lofar .glc files, in Lofar order and
    # de-duplicated — the search space for the analysis-sheet fallback below.
    lofar_dirs: list[Path] = []

    for i, link in enumerate(gram.glc_links, start=1):
        warnings: list[str] = []
        href = link.href
        glc_path = ""
        time_end = ""
        bandwidth = ""
        bandcentre = ""
        png_path = ""
        display_text = link.display_text
        resolved = resolve_glc_path(href, content_root, source_dir=source_dir)
        if resolved is None:
            warnings.append("GLC not found")
            glc_path = href
        else:
            glc_path = _rel_to_root(resolved, content_root)
            if resolved.parent not in lofar_dirs:
                lofar_dirs.append(resolved.parent)
            glc = parse_glc(resolved)
            warnings.extend(glc.warnings)
            bandwidth = glc.bandwidth
            bandcentre = glc.bandcentre
            if glc.image_filename:
                # The inner asset (.png/.jpg pre-rendered spectrogram, or
                # .wav for live render) sits next to the GLC on disk.
                # Resolve against the GLC's directory so the path is
                # image-root-relative and the generator can copy it
                # directly (see dita-topic-schema.md §1.2/§1.3 for how
                # the generator dispatches on the extension).
                png_path = resolve_asset_path(
                    glc.image_filename, content_root, source_dir=resolved.parent,
                )
                # Issue #148: the gram's time period (the GramFrame y-axis,
                # ``time_end``) is the image's scan-line count — its pixel
                # height — not the GLC's ``bottom_crop``. Only a pre-rendered
                # image carries this; a .wav-backed row is surfaced as a link
                # to its .glc (no GramFrame table) and keeps a blank time_end.
                if Path(png_path).suffix.lower() in GRAMFRAME_GLC_EXTENSIONS:
                    height = _image_pixel_height(png_path, content_root)
                    if height is not None:
                        time_end = str(height)

        row = {
            "publication": publication,
            "chapter": chapter or "",
            "target_doc": target_doc,
            "target_chapter": target_chapter,
            "gram_id": gram.gram_id,
            "vessel_name": gram.vessel_name,
            "topic_type": "glc",
            "sequence": str(i),
            "topic_filename": f"gram_{gram_num}.dita",
            "display_text": display_text,
            "link_href": href,
            "glc_path": glc_path,
            "time_end": time_end,
            "bandwidth": bandwidth,
            "bandcentre": bandcentre,
            "png_path": png_path,
            "target_ext": Path(png_path).suffix.lower(),
            "file_size": _png_file_size(png_path, content_root),
            "wav_treatment": "",
            "warnings": "",
        }
        # A resolved asset reference (image or .wav) whose file is not on disk
        # would dangle through the generator and only fail at publish; flag it
        # here so the operator can track it. An assetless GLC row (png_path "",
        # already carrying "GLC not found") is exempt.
        asset_missing = _asset_file_missing(png_path, content_root)
        # An *image* GLC-backed gram needs its time + frequency view fields for
        # GramFrame; flag (or, under --relaxed, default) any blank one so the
        # issue surfaces here rather than later in dedupe (issue #92). A .wav row
        # is exempt -- it links to its .glc, no GramFrame table -- as is an
        # assetless GLC row (target_ext "") which dangles.
        #
        # ``time_end`` is handled apart from the GLC-authored band fields: it is
        # the image's pixel height (issue #148), so a blank one is an *asset*
        # problem, not an author omission. When the image is missing on disk it
        # already dangles (ASSET_MISSING below, and re-resolves when the file is
        # dropped in), so — per the "missing assets dangle, they don't crash"
        # invariant — it is exempt from the fail-fast gate; only a
        # present-but-unreadable image earns a (non-fatal) warning here. The
        # fail-fast gate (glc_view_problems) therefore covers band fields only.
        if row["target_ext"] in GRAMFRAME_GLC_EXTENSIONS:
            for field_name in GLC_VIEW_FIELDS:
                if (row[field_name] or "").strip():
                    continue
                if relaxed:
                    row[field_name] = RELAXED_DEFAULT
                    warnings.append(
                        f"gram missing {field_name} — defaulted to "
                        f"{RELAXED_DEFAULT} (--relaxed)")
                elif field_name == "time_end":
                    if not asset_missing:
                        warnings.append(
                            "gram time period unknown — could not read image "
                            "height")
                else:
                    warnings.append(
                        f"gram missing {field_name} — GramFrame cannot "
                        f"render")
        if asset_missing:
            warnings.append(ASSET_MISSING_WARNING)
        row["warnings"] = ", ".join(warnings)
        rows.append(row)

    # Demon rows (issue #151) lead the gram: emitted here, once the Lofar folders
    # are known, and prepended to every return path so they precede the Lofar
    # rows in the CSV (the generator renders them first, after the analysis
    # sheet). A gram with no resolved Lofar folder gets none.
    demon_rows = demon_rows_for_gram(
        gram, lofar_dirs, publication, chapter, content_root, gram_num,
        target_doc, target_chapter, relaxed)

    analysis_warnings: list[str] = []
    analysis_png = gram.png_href or ""
    if not analysis_png:
        analysis_warnings.append("missing analysis PNG hyperlink")
        analysis_png_resolved = ""
    else:
        analysis_png_resolved = resolve_asset_path(analysis_png, content_root, source_dir)
        # Feature 007: an analysis sheet authored as a Word document
        # (.doc/.docx) is rendered to a same-stem .png sibling by
        # snapshot_analysis_docs.py ahead of extraction. Redirect the
        # analysis row's png_path to that sibling so the unchanged generator
        # embeds it inline (FR-004). When the rendered .png is absent (the
        # snapshotter hasn't run, or its render failed) keep the intended
        # .png href so the image dangles -- never a Word <xref> -- and record
        # a warning the author sees in Excel (FR-009/FR-010).
        if Path(analysis_png_resolved).suffix.lower() in (".doc", ".docx"):
            analysis_png_resolved = str(
                PurePosixPath(analysis_png_resolved).with_suffix(".png"))
            if not (content_root / analysis_png_resolved).is_file():
                analysis_warnings.append("analysis image not rendered")
    # A directly-referenced analysis image (.png/.jpg) whose file is absent —
    # typically a legacy href left pointing at an earlier build's folder. The
    # Word-sheet path above already records the more specific "analysis image
    # not rendered" (snapshot step not run) for the same missing .png, so that
    # case is excluded here. Otherwise:
    #   * a gram carrying a vessel_name is real content worth persevering with:
    #     look for an analysis sheet sitting in the gram's own Lofar folder(s)
    #     and repoint to it (a stale link whose real target moved). Found ->
    #     recover silently; not found -> flag ASSET_MISSING_WARNING as before.
    #   * a gram with no vessel_name is mangled legacy content we can't exploit:
    #     drop the analysis sheet entirely (no row, no warning), keeping the
    #     gram's resolved Lofar rows.
    if (analysis_png_resolved
            and "analysis image not rendered" not in analysis_warnings
            and _asset_file_missing(analysis_png_resolved, content_root)):
        if gram.vessel_name.strip():
            recovered = _recover_analysis_sheet(lofar_dirs, content_root)
            if recovered is not None:
                LOGGER.info(
                    "Gram %s (%s): analysis link %r not on disk; recovered "
                    "analysis sheet from Lofar folder -> %s",
                    gram.gram_id or "<no gram_id>", gram.vessel_name,
                    analysis_png, recovered,
                )
                analysis_png_resolved = recovered
            else:
                analysis_warnings.append(ASSET_MISSING_WARNING)
        else:
            LOGGER.info(
                "Gram %s: no vessel_name and analysis link %r not on disk — "
                "dropping the analysis sheet (mangled legacy, unrecoverable). "
                "The gram's Lofar rows are unaffected.",
                gram.gram_id or "<no gram_id>", analysis_png,
            )
            return demon_rows + rows
    rows.append({
        "publication": publication,
        "chapter": chapter or "",
        "target_doc": target_doc,
        "target_chapter": target_chapter,
        "gram_id": gram.gram_id,
        "vessel_name": gram.vessel_name,
        "topic_type": "analysis",
        "sequence": "1",
        "topic_filename": f"gram_{gram_num}.dita",
        "display_text": "",
        "link_href": "",
        "glc_path": "",
        "time_end": "",
        "bandwidth": "",
        "bandcentre": "",
        "png_path": analysis_png_resolved,
        "target_ext": Path(analysis_png_resolved).suffix.lower(),
        "file_size": _png_file_size(analysis_png_resolved, content_root),
        "wav_treatment": "",
        "warnings": ", ".join(analysis_warnings),
    })
    return demon_rows + rows


def glc_view_problems(rows: list[dict]) -> list[tuple[str, str, int]]:
    """Return ``(gram_id, field, line)`` for every GLC gram row missing a view field.

    Every *image* GLC-backed gram (a ``glc`` row whose inner asset is a
    renderable image) must carry ``bandwidth`` and ``bandcentre`` for GramFrame
    to render. ``.wav`` rows are exempt -- they surface as a link to the
    ``.glc``, never a GramFrame table, so their view fields are optional.
    Analysis rows and assetless (dangling) GLC rows are also exempt. Under
    ``--relaxed`` the image blanks are already filled with ``RELAXED_DEFAULT``
    upstream, so this returns empty (issue #92).

    ``time_end`` is **not** gated here: it is the image's pixel height (issue
    #148), so a blank one is a missing/unreadable-asset problem that dangles
    (flagged ASSET_MISSING and resolved by dropping the file in), not an author
    omission — hard-failing on it would violate the "missing assets dangle"
    invariant. Only the GLC-authored band fields fail-fast.

    ``line`` is the 1-based line the row occupies in the CSV this run writes
    (``write_csv`` emits the header then ``rows`` in order, so index ``i`` lands
    on line ``i + 2``), letting the abort point straight at the offending row.
    """
    problems: list[tuple[str, str, int]] = []
    for index, row in enumerate(rows):
        if row.get("topic_type", "") != "glc":
            continue
        if (row.get("target_ext", "") or "").lower() not in GRAMFRAME_GLC_EXTENSIONS:
            continue
        line_no = index + 2  # +1 header, +1 for 0-based -> 1-based
        for field_name in GLC_AUTHORED_VIEW_FIELDS:
            if not (row.get(field_name, "") or "").strip():
                problems.append((row.get("gram_id", ""), field_name, line_no))
    return problems


def missing_asset_problems(
    rows: list[dict], content_root: Path,
) -> list[tuple[str, str, str, int]]:
    """Return ``(publication, gram_id, png_path, line)`` for each asset-missing row.

    A row *expects* a target file when its ``png_path`` is non-empty (a resolved
    image or ``.wav`` reference). Rows that legitimately carry no asset — an
    assetless GLC link, or a gram with no analysis hyperlink — have a blank
    ``png_path`` and are exempt (the user's "not all rows require a target
    file"). A populated reference whose file is absent on disk would dangle
    through the generator and only fail at publish, so it is collected here.

    ``line`` mirrors ``glc_view_problems``: ``write_csv`` emits the header then
    ``rows`` in order, so index ``i`` lands on CSV line ``i + 2`` — the report
    can point straight at the offending row.
    """
    problems: list[tuple[str, str, str, int]] = []
    for index, row in enumerate(rows):
        png_path = (row.get("png_path", "") or "").strip()
        if not png_path:
            continue
        if _asset_file_missing(png_path, content_root):
            problems.append((
                row.get("publication", ""), row.get("gram_id", ""),
                png_path, index + 2,
            ))
    return problems


# -----------------------------------------------------------------------------
# CSV writer (R11)
# -----------------------------------------------------------------------------

def write_csv(rows: list[dict], out: Path) -> None:
    """Write rows to a UTF-8-with-BOM, CRLF, QUOTE_MINIMAL CSV."""
    out.parent.mkdir(parents=True, exist_ok=True)
    with out.open("w", encoding="utf-8-sig", newline="") as fh:
        writer = csv.DictWriter(
            fh, fieldnames=list(CSV_COLUMNS),
            quoting=csv.QUOTE_MINIMAL, lineterminator="\r\n",
        )
        writer.writeheader()
        for row in rows:
            writer.writerow({col: row.get(col, "") for col in CSV_COLUMNS})


# -----------------------------------------------------------------------------
# Orchestration
# -----------------------------------------------------------------------------

def main(argv: Iterable[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Extract PPTX content into the intermediate CSV")
    parser.add_argument("--input-root", required=True, type=Path, dest="input_root")
    parser.add_argument("--out", required=True, type=Path)
    parser.add_argument("--test-pattern", default=DEFAULT_TEST_PATTERN, dest="test_pattern")
    parser.add_argument("--final-pattern", default=DEFAULT_FINAL_PATTERN, dest="final_pattern")
    parser.add_argument("--joining-pattern", default=DEFAULT_JOINING_PATTERN, dest="joining_pattern")
    parser.add_argument(
        "--only", default=None, dest="only",
        help="Scope the walk to PPTXs whose path under --input-root starts with "
             "this exact folder name. The output CSV's path schema stays "
             "corpus-root-relative, so dedupe/generate can keep --image-root at "
             "--input-root unchanged. Useful for fast per-chapter debug "
             "iteration without re-pointing the downstream wrappers.",
    )
    parser.add_argument(
        "--exclude-tests", action="store_true", dest="exclude_tests",
        help="Skip the progress-test and final-assessment decks, emitting only "
             "the 'main' publication. Lets you build and review the main "
             "document from the full corpus without first carving the tests "
             "out of source\\. Independent of --only.",
    )
    parser.add_argument(
        "--strict-assets", action="store_true", dest="strict_assets",
        help="Promote the missing-target-asset check to a hard abort. By "
             "default a row that references an image/.wav whose file is absent "
             "on disk is flagged (a per-row warning plus an enumerated summary) "
             "but the CSV is still written, honouring the 'missing assets "
             "dangle' rule. Pass this during a focused cleanup pass to make the "
             "run fail (exit 1) until every referenced asset is present or its "
             "row dropped.",
    )
    parser.add_argument(
        "--relaxed", action="store_true", dest="relaxed",
        help="Dev/exploration only: instead of hard-failing when an image "
             "GLC-backed gram (pre-rendered .png/.jpg) is missing a "
             f"required view field ({', '.join(GLC_VIEW_FIELDS)}), substitute the "
             f"default '{RELAXED_DEFAULT}' so the rest of the toolchain can run "
             "against an incomplete corpus. GramFrame needs the real values, so "
             "fix the GLC and re-run without --relaxed before generating "
             "deliverable output.",
    )
    args = parser.parse_args(list(argv) if argv is not None else None)

    setup_logging(Path("extract.log"))

    if not args.input_root.is_dir():
        LOGGER.error("Input root does not exist or is not a directory: %s", args.input_root)
        return 1

    # Clear any prior output up-front. The input root is validated above (a
    # typo there must not destroy a good CSV), but past that point we commit
    # to producing fresh output: removing the target both verifies it isn't
    # locked (e.g. open in Excel) and guarantees a failed or zero-match run
    # can't leave a previous document's extract.csv behind for the downstream
    # dedupe/generate steps to silently consume.
    if args.out.exists():
        LOGGER.info("Removing existing output CSV %s", args.out)
        args.out.unlink()

    rows: list[dict] = []
    pptx_count = 0
    skipped_deleted_grams = 0
    warning_counter: Counter[str] = Counter()
    allocated: dict[str, int] = {}
    final_allocated: dict[str, int] = {}
    joining_allocated: dict[str, int] = {}

    try:
        for pptx in walk_pptxs(args.input_root, only_subdir=args.only):
            pptx_count += 1
            LOGGER.info("Processing PPTX %s", pptx)
            publication, chapter, chapter_slug = classify_publication(
                pptx, args.test_pattern, allocated,
                args.final_pattern, final_allocated,
                args.joining_pattern, joining_allocated,
            )
            if args.exclude_tests and publication != "main":
                LOGGER.info(
                    "Skipping %s deck (--exclude-tests): %s", publication, pptx.name)
                continue
            try:
                prs = Presentation(pptx)
            except Exception as exc:
                LOGGER.error("Cannot open PPTX %s: %s", pptx, exc)
                return 1
            deck_grams: list = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                if is_framing_slide(slide):
                    LOGGER.info("Skipping framing slide %d in %s", slide_num, pptx.name)
                    continue
                deck_grams.extend(extract_grams_from_slide(
                    slide, slide_num,
                    content_root=args.input_root,
                    source_dir=pptx.parent,
                ))
            # Feature 008/009: ``main`` is organised into four week folders.
            # A deck whose title carries a ``Week N`` token sends all its grams
            # to that week; a no-week ``main`` deck (e.g. Pub10, Legacy Pub 10)
            # has its grams sliced **evenly** across the four weeks (feature
            # 009), with the editable ``target_chapter`` filled per gram. Either
            # way ``main`` has no per-document folder segment (target_doc empty).
            # Non-main publications keep their existing per-document layout.
            target_doc = "" if publication == "main" else pptx.name
            target_chapters = deck_target_chapters(publication, chapter, len(deck_grams))
            for gram, target_chapter in zip(deck_grams, target_chapters):
                # A gram whose whole asset folder an author has deleted is a
                # removed gram, not a dangling asset: drop it entirely rather
                # than emit a topic of references that only fail at publish.
                # (A single deleted *file* with the folder still present is
                # left to the per-row missing-asset flag below.)
                if _gram_folder_missing(
                        gram, args.input_root, source_dir=pptx.parent):
                    LOGGER.info(
                        "Skipping gram %s in %s: its whole asset folder is "
                        "missing on disk — an author has deleted the gram "
                        "(dropped from the CSV, not dangled).",
                        gram.gram_id or "<no gram_id>", pptx.name,
                    )
                    skipped_deleted_grams += 1
                    continue
                gram_rows = gram_to_rows(
                    gram, publication, chapter, chapter_slug,
                    args.input_root, source_dir=pptx.parent,
                    target_doc=target_doc,
                    target_chapter=target_chapter,
                    relaxed=args.relaxed,
                )
                rows.extend(gram_rows)
                for r in gram_rows:
                    if r["warnings"]:
                        for w in r["warnings"].split(", "):
                            warning_counter[w] += 1
    except NotImplementedError as exc:
        LOGGER.error("Shape grouping stub reached: %s", exc)
        # Still write whatever rows were collected before the stub fired.
        write_csv(rows, args.out)
        return 1

    if args.only is not None and pptx_count == 0:
        LOGGER.warning(
            "--only=%r matched no PPTXs under %s. Check the folder name "
            "(case-insensitive match against the first path segment under "
            "--input-root) and that the directory contains *.pptx files.",
            args.only, args.input_root,
        )
        return 1

    write_csv(rows, args.out)

    # Missing target assets: a row that references an image/.wav whose file is
    # absent on disk would dangle silently through the generator and only
    # surface as a DITA-OT load error at publish (e.g. an analysis sheet whose
    # PNG never made it into source). Enumerate them here, at the extraction
    # boundary, so the operator can track and triage each one — fix the source
    # file, or drop the row — before publishing. Per the "missing assets
    # dangle, they don't crash" invariant this is a WARNING that still yields a
    # CSV (each offending row is also flagged in its warnings column);
    # --strict-assets promotes it to a hard abort for a focused cleanup pass.
    missing_assets = missing_asset_problems(rows, args.input_root)
    if missing_assets:
        detail = "; ".join(
            f"CSV line {line_no}: {pub}/gram_id={gram_id!r} -> {png_path}"
            for pub, gram_id, png_path, line_no in missing_assets)
        if args.strict_assets:
            LOGGER.error(
                "Aborting (--strict-assets): %d row(s) reference an asset file "
                "missing on disk. Fix the source file so it is present under "
                "%s, or drop the row, then re-run. Offending rows: %s",
                len(missing_assets), args.input_root, detail,
            )
            return 1
        LOGGER.warning(
            "%d row(s) reference an asset file missing on disk — each will "
            "dangle in the generated DITA and fail to load in Oxygen/DITA-OT, "
            "and is flagged with %r in its CSV warnings column. Fix the source "
            "file or drop the row (re-run with --strict-assets to hard-fail on "
            "these). Offending rows: %s",
            len(missing_assets), ASSET_MISSING_WARNING, detail,
        )

    # Issue #92: fail at extraction (not late in dedupe) when a GLC-backed gram
    # lacks the time + frequency view fields GramFrame requires. The CSV is
    # written first so its warnings column is inspectable; --relaxed has already
    # filled the blanks, so this scan only bites in strict mode.
    problems = glc_view_problems(rows)
    if problems:
        detail = "; ".join(
            f"CSV line {line_no}: gram_id={gram_id!r} missing {field_name}"
            for gram_id, field_name, line_no in problems)
        LOGGER.error(
            "Aborting: %d GLC gram view field(s) missing — GramFrame cannot "
            "render without them. Fix the GLC(s) so they carry a bandwidth and "
            "bandcentre (the time period is read from the image height), or "
            "re-run with --relaxed to substitute '%s' for development. "
            "Offending rows: %s",
            len(problems), RELAXED_DEFAULT, detail,
        )
        return 1

    if skipped_deleted_grams:
        LOGGER.warning(
            "Skipped %d gram(s) whose whole asset folder is missing on disk "
            "(author-deleted grams, dropped from the CSV). A single missing "
            "file with its folder still present is not dropped — it is flagged "
            "%r instead.",
            skipped_deleted_grams, ASSET_MISSING_WARNING,
        )

    distinct = ", ".join(f"{w}={c}" for w, c in sorted(warning_counter.items()))
    LOGGER.info(
        "Extraction summary: pptx=%d rows=%d skipped_deleted_grams=%d "
        "warnings=%d distinct=[%s]",
        pptx_count, len(rows), skipped_deleted_grams,
        sum(warning_counter.values()), distinct,
    )
    return 0


if __name__ == "__main__":
    rc = main()
    # Preserve CLI exit codes when invoked as a script, but stay silent
    # when invoked from an interactive REPL via runpy.run_path —
    # ``sys.exit`` would otherwise kill the interpreter and break the
    # up-arrow iteration loop. ``sys.ps1`` is only defined in
    # interactive sessions.
    if rc and not hasattr(sys, "ps1"):
        sys.exit(rc)

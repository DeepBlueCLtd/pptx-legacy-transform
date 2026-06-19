"""Synthetic instructor PPTX corpus generator (User Story 4, Phase 10 redesign).

Produces a multi-publication mock corpus structurally faithful to the real
source disk described in ``source/notes/reverse-spec.md`` — across three
families (Weeks, Progress Tests, Final Assessment, Pub10_Ed22B), each with
its own ``<Publication>.pptx`` plus a sibling ``<Publication> Files/`` tree
of ``Gram N/`` folders containing GLC + image / WAV / Analysis Sheet
assets.

Run as a script:

    python scripts/mock_pptx.py --out-root source

Defaults to writing into the repo's ``source/`` folder so the corpus is a
committed, reviewable deliverable (per the Phase 10 amendment in
``specs/001-pptx-dita-migration/tasks.md``).

The corpus uses a fixed RNG seed (``RANDOM_SEED``) so structural choices
(gram counts, Lofar counts, vessel/codename selection, analysis-sheet
type per gram) are reproducible. PPTX files themselves are not strictly
byte-identical across runs because python-pptx writes a small amount of
state into the OOXML zip; the *content* is deterministic.

All non-PPTX assets (``.glc``, ``.png``, ``.wav``, ``.docx``) are emitted
using the Python standard library only — no Pillow, no python-docx. The
GLC follows ``contracts/glc-schema.md`` exactly so the parser tests can
read every emitted file.
"""

from __future__ import annotations

import argparse
import io
import random
import shutil
import struct
import sys
import wave
import zipfile
import zlib
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from lxml import etree


# -----------------------------------------------------------------------------
# Determinism
# -----------------------------------------------------------------------------

RANDOM_SEED: int = 20260515


# -----------------------------------------------------------------------------
# Corpus shape (reverse-spec §1, §7)
# -----------------------------------------------------------------------------

@dataclass(frozen=True)
class FamilyParams:
    grams: int           # approximate gram count
    slides: int          # content slide count (welcome + exit are added on top)
    grams_per_slide: int  # nominal layout target


# Framing slide subtitle text — also used by the extractor to detect and
# skip welcome / exit slides during CSV generation.
FRAMING_SUBTITLE: str = "Instructor Version"
WELCOME_TITLE_PREFIX: str = "Welcome to "
EXIT_TITLE_PREFIX: str = "End of "


FAMILY_WEEK = FamilyParams(grams=35, slides=3, grams_per_slide=15)
FAMILY_TEST = FamilyParams(grams=30, slides=2, grams_per_slide=15)
FAMILY_FINAL = FamilyParams(grams=40, slides=3, grams_per_slide=15)
FAMILY_PUB10 = FamilyParams(grams=75, slides=5, grams_per_slide=15)


@dataclass(frozen=True)
class Publication:
    name: str            # folder + .pptx stem
    family: FamilyParams
    no_fr: bool = False  # rewrite "FR " prefix out of descriptors
    batched: bool = False  # Pub10 ten-gram folder batching


PUBLICATIONS: tuple[Publication, ...] = (
    Publication("Instructor Week 1 Grams", FAMILY_WEEK),
    Publication("Instructor Week 2 Grams", FAMILY_WEEK),
    Publication("Instructor Week 3 Grams", FAMILY_WEEK),
    Publication("Instructor Week 4 Grams_Updated", FAMILY_WEEK),
    Publication("Instructor Progress Test 1 Grams", FAMILY_TEST),
    Publication("Instructor Progress Test 2 Grams_Updated", FAMILY_TEST),
    Publication("Instructor Progress Test 3 Grams", FAMILY_TEST),
    Publication("Instructor Progress Test 3 Grams No FR", FAMILY_TEST, no_fr=True),
    Publication("Instructor Progress Test 4 Grams", FAMILY_TEST),
    Publication("Instructor Progress Final Assessment Grams", FAMILY_FINAL),
    Publication("Instructor Pub10_Ed22B_Updated", FAMILY_PUB10, batched=True),
)


# -----------------------------------------------------------------------------
# Vocabulary (reverse-spec §6)
# -----------------------------------------------------------------------------

# Vessel names / classes — Star Trek + Star Wars
VESSEL_NAMES: tuple[str, ...] = (
    # Star Trek
    "Enterprise", "Defiant", "Voyager", "Discovery", "Cerritos",
    "Reliant", "Excelsior", "Stargazer", "Equinox", "Prometheus",
    "Constitution-class", "Galaxy-class", "Intrepid-class", "Sovereign-class",
    "Nebula-class", "Akira-class", "Miranda-class", "Defiant-class",
    # Star Wars
    "Millennium Falcon", "Tantive IV", "Ghost", "Razor Crest",
    "Slave I", "Outrider", "Rebel One", "Devastator",
    "Star Destroyer", "X-wing", "Y-wing", "B-wing", "A-wing",
    "TIE Fighter", "TIE Bomber", "Lambda Shuttle", "Corellian Corvette",
    "Mon Calamari Cruiser", "Imperial-class", "Venator-class",
    "Acclamator-class", "Nebulon-B",
)

# Codenames — short fictional tokens reused across publications
CODENAMES: tuple[str, ...] = (
    "Tantive", "Defiant", "Tatooine", "Endor", "Hoth", "Yavin",
    "Dagobah", "Coruscant", "Naboo", "Bespin", "Mustafar",
    "Romulus", "Vulcan", "Bajor", "Cardassia", "Kronos",
    "Risa", "Trill", "Betazed", "Kobol", "Caprica",
    "Dank", "Gandalf",  # deliberate cosmetic noise (cf. reverse-spec §7 examples)
)


# -----------------------------------------------------------------------------
# Slide geometry
# -----------------------------------------------------------------------------

SLIDE_WIDTH_IN: float = 13.33
SLIDE_HEIGHT_IN: float = 7.5
TITLE_BAR_HEIGHT_IN: float = 0.6
GRID_TOP_MARGIN_IN: float = 0.8
GRID_LEFT_MARGIN_IN: float = 0.4
GRID_ROWS: int = 3
GRID_COLS: int = 5
TILE_TITLE_HEIGHT_IN: float = 0.4
TILE_LOFAR_LINE_HEIGHT_IN: float = 0.25
CELL_GAP_IN: float = 0.1


# -----------------------------------------------------------------------------
# Hyperlink helpers (R5; both shape-level and text-run forms exercised)
# -----------------------------------------------------------------------------

HYPERLINK_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink"
)


def add_shape_level_hyperlink(shape, target: str) -> None:
    """Attach a shape-level click action via direct lxml manipulation.

    The ``a:hlinkClick`` element belongs inside ``p:cNvPr`` (NonVisualDrawingProps)
    per ECMA-376. PowerPoint ignores it when placed under ``p:nvPr``.
    """
    rel_id = shape.part.relate_to(target, HYPERLINK_REL_TYPE, is_external=True)
    nv_sp_pr = shape._element.find(qn("p:nvSpPr"))
    c_nv_pr = nv_sp_pr.find(qn("p:cNvPr"))
    hlink = etree.SubElement(c_nv_pr, qn("a:hlinkClick"))
    hlink.set(qn("r:id"), rel_id)


def add_text_run_hyperlink(run, target: str) -> None:
    """Attach a text-run hyperlink. python-pptx handles relative paths fine."""
    run.hyperlink.address = target


# -----------------------------------------------------------------------------
# Asset emitters (stdlib only; reverse-spec §5)
# -----------------------------------------------------------------------------

def emit_glc(
    path: Path, *, image_filename: str, time_end: int,
    bandwidth: int, bandcentre: int,
) -> None:
    """Write a GAPS_Lite_configuration XML file per contracts/glc-schema.md."""
    root = ET.Element("GAPS_Lite_configuration")
    ds = ET.SubElement(root, "data_source")
    ET.SubElement(ds, "filename").text = image_filename
    crops = ET.SubElement(ds, "bitmap_crop_values")
    ET.SubElement(crops, "top_crop").text = "0"
    ET.SubElement(crops, "bottom_crop").text = str(time_end)
    pb = ET.SubElement(root, "playback")
    ET.SubElement(pb, "time_offset").text = "0"
    settings = ET.SubElement(root, "settings")
    lofar = ET.SubElement(settings, "lofar")
    ET.SubElement(lofar, "bandwidth").text = str(bandwidth)
    ET.SubElement(lofar, "bandcentre").text = str(bandcentre)
    tree = ET.ElementTree(root)
    path.parent.mkdir(parents=True, exist_ok=True)
    tree.write(path, encoding="utf-8", xml_declaration=True)


def _png_chunk(tag: bytes, data: bytes) -> bytes:
    return struct.pack(">I", len(data)) + tag + data + struct.pack(
        ">I", zlib.crc32(tag + data) & 0xFFFFFFFF
    )


def emit_png(path: Path, *, width: int = 8, height: int = 8, shade: int = 200) -> None:
    """Write a tiny valid grayscale PNG using only stdlib (no Pillow).

    Used as the deterministic fallback when ``MOCK_SPECTROGRAMS_DIR`` is
    missing — the synthetic block is enough to keep the pipeline tests
    self-contained, but it renders as a near-invisible square in the
    browser. Prefer :func:`emit_spectrogram` for any output that will be
    inspected visually (gh-pages demo, screenshots, etc.).
    """
    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = struct.pack(">IIBBBBB", width, height, 8, 0, 0, 0, 0)  # 8-bit grayscale
    raw = b""
    for _ in range(height):
        raw += b"\x00" + bytes([shade] * width)  # filter byte + pixels
    idat = zlib.compress(raw, 9)
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("wb") as fh:
        fh.write(sig)
        fh.write(_png_chunk(b"IHDR", ihdr))
        fh.write(_png_chunk(b"IDAT", idat))
        fh.write(_png_chunk(b"IEND", b""))


MOCK_SPECTROGRAMS_DIR = Path(__file__).resolve().parent / "mock_pptx_data"
SPECTROGRAM_VARIANT_COUNT = 4


def _pick_variant(stem: str) -> int:
    """Map a filename stem onto one of the ``SPECTROGRAM_VARIANT_COUNT``
    pre-baked spectrogram variants. Deterministic so two consecutive
    mock-corpus runs produce byte-identical assets (R9 idempotency)."""
    return zlib.adler32(stem.encode("utf-8")) % SPECTROGRAM_VARIANT_COUNT


def emit_spectrogram(path: Path) -> None:
    """Copy one of the pre-baked spectrogram variants to ``path``.

    The variant is chosen by ``_pick_variant`` so every gram-and-Lofar
    pair gets a stable, visually distinct image across runs. Falls back
    to a synthetic placeholder when the asset dir is absent (bare-clone
    tests, packaged distributions without the data dir, etc.) so the
    pipeline still runs end-to-end.
    """
    path.parent.mkdir(parents=True, exist_ok=True)
    variant = _pick_variant(path.stem)
    source = MOCK_SPECTROGRAMS_DIR / f"spectrogram-{variant}.png"
    if source.is_file():
        shutil.copy2(source, path)
        return
    emit_png(path, width=8, height=8, shade=120 + variant * 30)


def emit_analysis_sheet_image(path: Path) -> None:
    """Copy the pre-baked Analysis Sheet image to ``path``.

    Real-world Analysis Sheets are screenshots of an MS Word measurement
    table (harmonics, base frequencies, speed, etc.) — not spectrograms.
    Falls back to ``emit_spectrogram`` when the dedicated asset is
    missing so older mock-data dirs keep working.
    """
    path.parent.mkdir(parents=True, exist_ok=True)
    source = MOCK_SPECTROGRAMS_DIR / "analysis-sheet.png"
    if source.is_file():
        shutil.copy2(source, path)
        return
    emit_spectrogram(path)


def emit_wav(path: Path, *, duration_s: float = 0.1, framerate: int = 8000) -> None:
    """Write a short silence WAV using stdlib ``wave``."""
    path.parent.mkdir(parents=True, exist_ok=True)
    frames = int(duration_s * framerate)
    with wave.open(str(path), "wb") as wf:
        wf.setnchannels(1)
        wf.setsampwidth(1)
        wf.setframerate(framerate)
        wf.writeframes(b"\x80" * frames)  # silence (mid-level for 8-bit unsigned)


_DOCX_CONTENT_TYPES = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
<Default Extension="xml" ContentType="application/xml"/>
<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""

_DOCX_ROOT_RELS = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""

_DOCX_BODY_TEMPLATE = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
<w:body>
<w:p><w:r><w:t xml:space="preserve">{title}</w:t></w:r></w:p>
<w:p><w:r><w:t xml:space="preserve">Bearing: __</w:t></w:r></w:p>
<w:p><w:r><w:t xml:space="preserve">Frequency: __</w:t></w:r></w:p>
<w:p><w:r><w:t xml:space="preserve">Identification: __</w:t></w:r></w:p>
<w:p><w:r><w:t xml:space="preserve">Notes: __</w:t></w:r></w:p>
</w:body>
</w:document>"""


def emit_docx(path: Path, *, title: str) -> None:
    """Write a minimal valid .docx (Word document) using stdlib zipfile."""
    path.parent.mkdir(parents=True, exist_ok=True)
    body = _DOCX_BODY_TEMPLATE.format(title=_escape_xml(title))
    # Use a fixed timestamp inside the zip for deterministic output.
    fixed_date = (1980, 1, 1, 0, 0, 0)
    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zf:
        for name, content in (
            ("[Content_Types].xml", _DOCX_CONTENT_TYPES),
            ("_rels/.rels", _DOCX_ROOT_RELS),
            ("word/document.xml", body),
        ):
            info = zipfile.ZipInfo(filename=name, date_time=fixed_date)
            info.compress_type = zipfile.ZIP_DEFLATED
            zf.writestr(info, content)


def emit_doc(path: Path, *, title: str) -> None:
    """Write deterministic placeholder bytes for a legacy ``.doc`` sheet.

    The mock is not a renderer and nothing downstream parses the ``.doc``
    (``snapshot_analysis_docs.py`` only hands it to LibreOffice), so a
    fixed, byte-stable placeholder is sufficient for the test corpus. The
    rendered sibling ``.png`` is emitted alongside by ``_emit_analysis_sheet``
    so the full pipeline exercises the doc -> inline-image path without
    LibreOffice (feature 007 T013).
    """
    path.parent.mkdir(parents=True, exist_ok=True)
    placeholder = f"[legacy .doc analysis sheet placeholder] {title}\n"
    path.write_bytes(placeholder.encode("utf-8"))


def _escape_xml(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


# -----------------------------------------------------------------------------
# Per-publication generation
# -----------------------------------------------------------------------------

@dataclass
class GramSpec:
    gram_num: int            # surviving original gram number (gaps allowed)
    folder_name: str         # on-disk folder name within Files/
    folder_path: Path        # absolute folder path on disk
    descriptor: str          # "Gram N: <instructor detail>"
    rel_analysis: str        # relative href from PPTX to Analysis Sheet
    analysis_kind: str       # "doc", "docx", or "png"
    lofars: list["LofarSpec"]


@dataclass
class LofarSpec:
    label: str               # "Lofar 1", etc.
    rel_glc: str             # relative href from PPTX to .glc
    rel_media: str           # relative href to .png or .wav referenced by GLC
    media_kind: str          # "png" or "wav"


# -----------------------------------------------------------------------------
# Deliberate cross-publication duplicates (reverse-spec §7; corpus dedup demo)
# -----------------------------------------------------------------------------

@dataclass(frozen=True)
class SharedLofar:
    label: str
    media_kind: str  # "png" or "wav"
    time_end: int
    bandwidth: int
    bandcentre: int


@dataclass(frozen=True)
class SharedGram:
    """A gram emitted byte-identically into two publications.

    Same gram_num, same lofar layout, same time_end / bandwidth / bandcentre values, same
    analysis-sheet kind, so every emitted asset (GLC, PNG, WAV, Analysis.png)
    is byte-for-byte identical across the two publications. The PPTX
    descriptor may differ between publications — that's the realistic case
    the duplicate-detection workflow has to catch (the file_size column
    surfaces it even when names drift).
    """
    gram_num: int
    descriptor_week1: str   # PPTX descriptor used in Instructor Week 1 Grams
    descriptor_test1: str   # PPTX descriptor used in Instructor Progress Test 1 Grams
    lofars: tuple[SharedLofar, ...]


SHARED_GRAMS_WEEK1_TEST1: tuple[SharedGram, ...] = (
    # Identical descriptors in both publications — pure dupe (matches by
    # both name and file_size).
    SharedGram(
        gram_num=3,
        descriptor_week1="Gram 3: FR Constitution-class, Category 2, Endor",
        descriptor_test1="Gram 3: FR Constitution-class, Category 2, Endor",
        lofars=(
            SharedLofar(label="Lofar 1", media_kind="png", time_end=300, bandwidth=400, bandcentre=200),
            SharedLofar(label="Lofar 2", media_kind="png", time_end=240, bandwidth=200, bandcentre=300),
        ),
    ),
    # Identical descriptors in both publications, but a wav-rendered lofar.
    SharedGram(
        gram_num=7,
        descriptor_week1="Gram 7: FR Defiant, Category 1, Vulcan",
        descriptor_test1="Gram 7: FR Defiant, Category 1, Vulcan",
        lofars=(
            SharedLofar(label="Lofar 1", media_kind="wav", time_end=271, bandwidth=100, bandcentre=150),
        ),
    ),
    # Slightly different descriptor across publications — same assets, drifted
    # name. This is the case the file_size column has to catch: gram_id and
    # vessel_name would not match by string compare, but every asset still has
    # an identical byte count to its twin in the other publication.
    SharedGram(
        gram_num=11,
        descriptor_week1="Gram 11: FR Voyager, Category 3, Risa",
        descriptor_test1="Gram 11: FR Voyager, Category 3, Risa (mark II)",
        lofars=(
            SharedLofar(label="Lofar 1", media_kind="png", time_end=180, bandwidth=800, bandcentre=400),
            SharedLofar(label="Lofar 2", media_kind="png", time_end=240, bandwidth=400, bandcentre=200),
            SharedLofar(label="Lofar 3", media_kind="png", time_end=360, bandwidth=200, bandcentre=300),
        ),
    ),
)

SHARED_PUB_NAMES: frozenset[str] = frozenset((
    "Instructor Week 1 Grams",
    "Instructor Progress Test 1 Grams",
))


def _gram_numbers_with_gaps(target: int, rng: random.Random) -> list[int]:
    """Return ``target`` numbers from a sequence with ~5-10% gaps removed.

    Models the reverse-spec's "occasional gaps from simulated edits" rule:
    pick a slightly larger pool, then drop a handful of integers.
    """
    pool_size = int(target * 1.15) + 2
    pool = list(range(1, pool_size + 1))
    drop_count = pool_size - target
    drops = rng.sample(pool[1:-1], drop_count) if drop_count > 0 else []
    return [n for n in pool if n not in drops][:target]


def _pick_descriptor(
    gram_num: int, rng: random.Random, vessel_pool: list[str], no_fr: bool
) -> str:
    """Compose a `"Gram N: <instructor detail>"` descriptor with deliberate format variance."""
    vessel = rng.choice(vessel_pool)
    codename = rng.choice(CODENAMES)
    category = rng.randint(1, 4)
    style = rng.choice(("fielded", "fielded", "fielded", "sentence"))
    if style == "fielded":
        prefix = "" if no_fr else "FR "
        detail = f"{prefix}{vessel}, Category {category}, {codename}"
    else:
        detail = f"{vessel} contact bearing {rng.randint(0, 359):03d}, codename {codename}"
    return f"Gram {gram_num}: {detail}"


def _build_lofars(
    gram_folder: Path,
    rel_folder_from_pptx: str,
    rng: random.Random,
) -> list[LofarSpec]:
    """Decide Lofar count (1-4) and emit each Lofar's GLC + media asset."""
    count = rng.randint(1, 4)
    lofars: list[LofarSpec] = []
    for i in range(1, count + 1):
        # ~80% PNG, 20% WAV per reverse-spec §7.
        media_kind = "wav" if rng.random() < 0.2 else "png"
        # Filename variants: suffix tokens for realism (reverse-spec §7).
        suffix = rng.choice(("", "_a", "_b", " I", " ABC", " Loop 1", " Loop 2"))
        media_stem = f"Lofar {i}{suffix}".strip()
        glc_path = gram_folder / f"{media_stem}.glc"
        media_path = gram_folder / f"{media_stem}.{media_kind}"
        time_end = rng.choice((180, 240, 271, 300, 360))
        bandwidth = rng.choice((100, 200, 400, 800))
        # Frequency band = bandwidth + bandcentre (issue #87). Mix centred bands
        # (centre == bandwidth/2 → starts at 0) with off-centre bands so the
        # synthetic corpus exercises the non-zero freq_start path.
        bandcentre = rng.choice((bandwidth // 2, bandwidth))
        # GLC references the sibling media file by basename (matches real-corpus convention).
        emit_glc(glc_path, image_filename=media_path.name,
                 time_end=time_end, bandwidth=bandwidth, bandcentre=bandcentre)
        if media_kind == "png":
            emit_spectrogram(media_path)
        else:
            emit_wav(media_path)
        lofars.append(LofarSpec(
            label=f"Lofar {i}",
            rel_glc=f"{rel_folder_from_pptx}/{glc_path.name}",
            rel_media=f"{rel_folder_from_pptx}/{media_path.name}",
            media_kind=media_kind,
        ))
    return lofars


def _emit_analysis_sheet(
    gram_folder: Path,
    rel_folder_from_pptx: str,
    kind: str,
    title: str,
) -> str:
    """Emit the analysis sheet and return its relative href from the PPTX."""
    if kind == "docx":
        name = "Analysis Sheet.docx"
        emit_docx(gram_folder / name, title=title)
    elif kind == "doc":
        name = "Analysis Sheet.doc"
        emit_doc(gram_folder / name, title=title)
        # Ship the rendered sibling .png (same stem) so the doc -> inline-image
        # path runs end to end without LibreOffice (feature 007 T013).
        emit_analysis_sheet_image(gram_folder / "Analysis Sheet.png")
    else:
        name = "Analysis.png"
        emit_analysis_sheet_image(gram_folder / name)
    return f"{rel_folder_from_pptx}/{name}"


def build_gram_specs(pub: Publication, files_dir: Path, rng: random.Random) -> list[GramSpec]:
    """Plan the gram folders + assets for one publication and write the files."""
    gram_nums = _gram_numbers_with_gaps(pub.family.grams, rng)

    # Cross-publication vessel reuse: keep a stable shuffled pool per publication
    # so the same vessel reappears across grams (and, by RNG construction, across
    # publications since we draw from the same VESSEL_NAMES tuple).
    vessel_pool = list(VESSEL_NAMES)
    rng.shuffle(vessel_pool)

    specs: list[GramSpec] = []
    for gram_num in gram_nums:
        folder_name, rel_folder = _gram_folder_layout(pub, gram_num, len(specs), rng)
        folder_path = files_dir / folder_name if not pub.batched else files_dir / rel_folder.split("/", 1)[0] / rel_folder.split("/", 1)[1]
        # For the batched case rel_folder already encodes "<batch>/<gram>".
        # For flat case the rel_folder is just the gram folder name.
        if pub.batched:
            batch_dir, leaf = rel_folder.split("/", 1)
            folder_path = files_dir / batch_dir / leaf
        else:
            folder_path = files_dir / rel_folder
        folder_path.mkdir(parents=True, exist_ok=True)

        descriptor = _pick_descriptor(gram_num, rng, vessel_pool, pub.no_fr)
        # Hyperlink hrefs are relative to the PPTX (sibling to "<Publication> Files/").
        files_dir_name = files_dir.name
        rel_folder_from_pptx = f"{files_dir_name}/{rel_folder}"

        analysis_kind = rng.choice(("doc", "docx", "png"))
        rel_analysis = _emit_analysis_sheet(
            folder_path, rel_folder_from_pptx, analysis_kind, title=descriptor)
        lofars = _build_lofars(folder_path, rel_folder_from_pptx, rng)

        specs.append(GramSpec(
            gram_num=gram_num,
            folder_name=folder_path.name,
            folder_path=folder_path,
            descriptor=descriptor,
            rel_analysis=rel_analysis,
            analysis_kind=analysis_kind,
            lofars=lofars,
        ))
    return specs


def _gram_folder_layout(
    pub: Publication, gram_num: int, index: int, rng: random.Random
) -> tuple[str, str]:
    """Return ``(folder_name, rel_path_under_Files)`` for one gram."""
    if not pub.batched:
        name = f"Gram {gram_num}"
        return name, name
    # Pub10_Ed22B: batches of 10. ``index`` is 0-based across surviving grams,
    # which maps to batch boundaries cleanly because the SME prunes/edits
    # surviving grams but the folders are still grouped by original tens.
    batch_lo = ((gram_num - 1) // 10) * 10 + 1
    batch_hi = batch_lo + 9
    batch = f"Pub 10_Ed 2_({batch_lo}-{batch_hi})"
    # Mild folder-name variance reverse-spec §7 calls out (zero-pad / codename suffix).
    style = rng.choice(("plain", "plain", "padded", "codename"))
    if style == "padded":
        leaf = f"Gram_{gram_num:02d}"
    elif style == "codename":
        leaf = f"Gram_{gram_num} {rng.choice(CODENAMES)}"
    else:
        leaf = f"Gram_{gram_num}"
    return leaf, f"{batch}/{leaf}"


# -----------------------------------------------------------------------------
# PPTX construction
# -----------------------------------------------------------------------------

def _add_title_bar(slide, prs: Presentation, *, title: str) -> None:
    """Add a coloured title bar with placeholder text (no real org imagery)."""
    width = prs.slide_width
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, width, Inches(TITLE_BAR_HEIGHT_IN)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    bar.line.fill.background()
    tf = bar.text_frame
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _add_framing_slide(prs: Presentation, *, title: str) -> None:
    """Add a welcome/exit slide: centred title + 'Instructor Version' subtitle.

    No gram content, no hyperlinks. The extractor identifies these slides
    by the title prefix and skips them when building the CSV.
    """
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    title_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(2.6),
        Inches(SLIDE_WIDTH_IN - 2.0), Inches(1.4),
    )
    ttf = title_box.text_frame
    ttf.word_wrap = True
    tp = ttf.paragraphs[0]
    tp.alignment = PP_ALIGN.CENTER
    trun = tp.add_run()
    trun.text = title
    trun.font.size = Pt(40)
    trun.font.bold = True
    trun.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    sub_box = slide.shapes.add_textbox(
        Inches(1.0), Inches(4.2),
        Inches(SLIDE_WIDTH_IN - 2.0), Inches(0.8),
    )
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    srun = sp.add_run()
    srun.text = FRAMING_SUBTITLE
    srun.font.size = Pt(24)
    srun.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)


def _add_gram_tile(
    slide, *, left_in: float, top_in: float, width_in: float, height_in: float,
    spec: GramSpec,
) -> None:
    """Place one gram tile: rounded rectangle + Lofar text labels beneath."""
    # Rounded rectangle holding the descriptor; shape-level hyperlink to Analysis.
    title_h = TILE_TITLE_HEIGHT_IN
    rect = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(title_h),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xE8, 0xEE, 0xF7)
    rect.line.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    tf = rect.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = spec.descriptor
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)
    add_shape_level_hyperlink(rect, spec.rel_analysis)

    # Lofar text box beneath the rectangle.
    lofar_top = top_in + title_h + 0.02
    lofar_h = max(0.15, height_in - title_h - 0.05)
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(lofar_top),
        Inches(width_in), Inches(lofar_h),
    )
    btf = box.text_frame
    btf.word_wrap = True
    for i, lofar in enumerate(spec.lofars):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        run = para.add_run()
        run.text = lofar.label
        run.font.size = Pt(8)
        add_text_run_hyperlink(run, lofar.rel_glc)


def _chunk(seq: list, size: int) -> Iterable[list]:
    for i in range(0, len(seq), size):
        yield seq[i:i + size]


def build_publication_pptx(pub: Publication, out_pptx: Path, specs: list[GramSpec]) -> None:
    """Build one publication's PPTX, embedding tiles + hyperlinks."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)

    slide_chunks = list(_chunk(specs, pub.family.grams_per_slide))
    total_pages = len(slide_chunks)
    blank_layout = prs.slide_layouts[6]  # blank

    _add_framing_slide(prs, title=f"{WELCOME_TITLE_PREFIX}{pub.name}")

    for page_num, chunk in enumerate(slide_chunks, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _add_title_bar(
            slide, prs,
            title=f"{pub.name} — Page {page_num} of {total_pages}",
        )

        # Grid placement.
        usable_w = SLIDE_WIDTH_IN - 2 * GRID_LEFT_MARGIN_IN
        usable_h = SLIDE_HEIGHT_IN - GRID_TOP_MARGIN_IN - 0.3
        cell_w = usable_w / GRID_COLS
        cell_h = usable_h / GRID_ROWS
        for idx, spec in enumerate(chunk):
            row = idx // GRID_COLS
            col = idx % GRID_COLS
            left = GRID_LEFT_MARGIN_IN + col * cell_w
            top = GRID_TOP_MARGIN_IN + row * cell_h
            _add_gram_tile(
                slide,
                left_in=left, top_in=top,
                width_in=cell_w - CELL_GAP_IN,
                height_in=cell_h - CELL_GAP_IN,
                spec=spec,
            )

    _add_framing_slide(prs, title=f"{EXIT_TITLE_PREFIX}{pub.name}")

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_pptx)


def _emit_shared_gram(
    files_dir: Path, shared: SharedGram, descriptor: str
) -> GramSpec:
    """Emit one duplicate-across-publications gram into ``files_dir``.

    Always uses the flat (non-batched) layout because the two participating
    publications (Week 1, Progress Test 1) are non-batched. Lofar media
    filenames match their labels exactly (no random suffix) so the same
    ``_pick_variant`` stem is used in both publications and the resulting
    PNG bytes are identical.
    """
    folder_name = f"Gram {shared.gram_num}"
    folder_path = files_dir / folder_name
    folder_path.mkdir(parents=True, exist_ok=True)

    rel_folder_from_pptx = f"{files_dir.name}/{folder_name}"

    lofars: list[LofarSpec] = []
    for lf in shared.lofars:
        media_stem = lf.label  # e.g. "Lofar 1" — no suffix; ensures byte-identical PNGs.
        glc_path = folder_path / f"{media_stem}.glc"
        media_path = folder_path / f"{media_stem}.{lf.media_kind}"
        emit_glc(glc_path, image_filename=media_path.name,
                 time_end=lf.time_end, bandwidth=lf.bandwidth, bandcentre=lf.bandcentre)
        if lf.media_kind == "png":
            emit_spectrogram(media_path)
        else:
            emit_wav(media_path)
        lofars.append(LofarSpec(
            label=lf.label,
            rel_glc=f"{rel_folder_from_pptx}/{glc_path.name}",
            rel_media=f"{rel_folder_from_pptx}/{media_path.name}",
            media_kind=lf.media_kind,
        ))

    # Always png-form analysis sheet so the bytes match across publications
    # (the docx-form's bytes depend on the descriptor, which may differ).
    rel_analysis = _emit_analysis_sheet(
        folder_path, rel_folder_from_pptx, "png", title=descriptor,
    )

    return GramSpec(
        gram_num=shared.gram_num,
        folder_name=folder_name,
        folder_path=folder_path,
        descriptor=descriptor,
        rel_analysis=rel_analysis,
        analysis_kind="png",
        lofars=lofars,
    )


def _inject_shared_duplicates(
    pub: Publication, files_dir: Path, specs: list[GramSpec]
) -> list[GramSpec]:
    """Replace any conflicting auto-generated specs with the shared duplicates.

    Only applies to Week 1 Grams and Progress Test 1 Grams; every other
    publication is returned unchanged. The shared specs are placed at the
    front of the list so they land on the first content slide.
    """
    if pub.name not in SHARED_PUB_NAMES:
        return specs
    is_week1 = pub.name == "Instructor Week 1 Grams"

    shared_nums = {s.gram_num for s in SHARED_GRAMS_WEEK1_TEST1}
    kept: list[GramSpec] = []
    for spec in specs:
        if spec.gram_num in shared_nums:
            if spec.folder_path.exists():
                shutil.rmtree(spec.folder_path)
            continue
        kept.append(spec)

    shared_specs: list[GramSpec] = []
    for shared in SHARED_GRAMS_WEEK1_TEST1:
        descriptor = shared.descriptor_week1 if is_week1 else shared.descriptor_test1
        shared_specs.append(_emit_shared_gram(files_dir, shared, descriptor))

    return shared_specs + kept


def build_publication(pub: Publication, out_root: Path, rng: random.Random) -> None:
    """Write one publication: <Name>/<Name>.pptx + <Name> Files/Gram N/..."""
    pub_dir = out_root / pub.name
    files_dir = pub_dir / f"{pub.name} Files"
    files_dir.mkdir(parents=True, exist_ok=True)
    specs = build_gram_specs(pub, files_dir, rng)
    specs = _inject_shared_duplicates(pub, files_dir, specs)
    out_pptx = pub_dir / f"{pub.name}.pptx"
    build_publication_pptx(pub, out_pptx, specs)


# -----------------------------------------------------------------------------
# Orchestration
# -----------------------------------------------------------------------------

def main(argv: Iterable[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Generate the synthetic instructor PPTX corpus")
    parser.add_argument(
        "--out-root", required=True, type=Path,
        help="Output directory; each publication becomes a subfolder",
    )
    args = parser.parse_args(list(argv) if argv is not None else None)

    out_root: Path = args.out_root
    out_root.mkdir(parents=True, exist_ok=True)

    rng = random.Random(RANDOM_SEED)

    for pub in PUBLICATIONS:
        # Use a publication-scoped sub-RNG so adding a publication later doesn't
        # disturb the earlier ones' generated content.
        sub_seed = rng.randrange(0, 2**31 - 1)
        sub_rng = random.Random(sub_seed)
        build_publication(pub, out_root, sub_rng)
        print(f"Built {pub.name}")

    print(f"Wrote {len(PUBLICATIONS)} publications under {out_root}")
    return 0


if __name__ == "__main__":
    rc = main()
    # Preserve CLI exit codes when invoked as a script, but stay silent
    # when invoked from an interactive REPL via runpy.run_path —
    # ``sys.exit`` would otherwise kill the interpreter and break the
    # up-arrow iteration loop. ``sys.ps1`` is only defined in
    # interactive sessions.
    if rc and not hasattr(sys, "ps1"):
        sys.exit(rc)

"""Structural-report producer for an instructor PPTX (User Story 3).

Default output is gram-focused — for each non-framing slide, the grams
the extractor would lift, with their analysis-sheet and .glc hrefs.
This is the view to use when verifying that a deck conforms to the
expected legacy structure.

``--verbose`` appends two forensic sections: a per-shape dump (every
shape's position, text, and hyperlinks) and a raw hyperlink-target
grouping (including the vestigial absolute ``file:///`` overlay links
that the default view filters out). Use it when a deck looks broken.

Sections:
  1. Summary: counts and anomaly flags.
  2. Per-gram (default): grams the extractor will lift.
  3. Per-shape (verbose only): every shape, position, text, hyperlinks.
  4. Raw hyperlinks (verbose only): every hyperlink target, grouped by
     extension and tagged shape-level / text-run.

Run as:

    python scripts/introspect_pptx.py --input PATH [--out PATH] [--slides N,M,...] [--verbose]
"""

from __future__ import annotations

import argparse
import logging
import sys
from collections import defaultdict
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.oxml.ns import qn

from extract_to_csv import (
    GramPlaceholder,
    _slide_diagram_hyperlinks,
    extract_grams_from_slide,
    is_framing_slide,
)


EXPECTED_SHAPES_PER_CONTENT_SLIDE: int = 30  # 15 titles + 15 link boxes
SHAPE_DEVIATION_TOLERANCE: int = 5
TEXT_TRUNCATION: int = 80

LOGGER = logging.getLogger(__name__)


def setup_logging(log_path: Path) -> None:
    """Configure dual stdout + per-stage-file logging per R10."""
    root = logging.getLogger()
    for handler in list(root.handlers):
        root.removeHandler(handler)
    formatter = logging.Formatter("%(asctime)s %(levelname)s %(name)s: %(message)s")
    stream = logging.StreamHandler(sys.stdout)
    stream.setLevel(logging.INFO)
    stream.setFormatter(formatter)
    file_handler = logging.FileHandler(log_path, mode="w", encoding="utf-8")
    file_handler.setLevel(logging.DEBUG)
    file_handler.setFormatter(formatter)
    root.setLevel(logging.DEBUG)
    root.addHandler(stream)
    root.addHandler(file_handler)


@dataclass
class RunRecord:
    text: str
    hyperlink: str | None


@dataclass
class ShapeRecord:
    slide_number: int
    index: int
    name: str
    shape_type: str
    left_in: float
    top_in: float
    width_in: float
    height_in: float
    text: str
    shape_hyperlink: str | None
    runs: list[RunRecord] = field(default_factory=list)


def extract_run_hyperlink(run) -> tuple[str | None, str]:
    """Return ``(target, "text-run")`` or ``(None, "")`` for a text run."""
    try:
        target = run.hyperlink.address
    except Exception:
        return (None, "")
    if not target:
        return (None, "")
    return (target, "text-run")


def extract_shape_hyperlink(shape) -> tuple[str | None, str]:
    """Return ``(target, "shape-level")`` or ``(None, "")`` for a shape click action.

    Searches every descendant ``p:cNvPr`` so the lookup works regardless
    of the shape's outer XML wrapper — ``p:sp`` (autoshape/textbox),
    ``p:pic`` (picture), ``p:cxnSp`` (connector), or ``p:graphicFrame``
    (chart/table/etc.). Pictures with shape-level clicks were
    previously invisible to this walker.
    """
    try:
        for c_nv_pr in shape._element.iter(qn("p:cNvPr")):
            hlink = c_nv_pr.find(qn("a:hlinkClick"))
            if hlink is None:
                continue
            rel_id = hlink.get(qn("r:id"))
            if not rel_id:
                continue
            target = shape.part.rels[rel_id].target_ref
            if not target:
                continue
            return (target, "shape-level")
        return (None, "")
    except Exception:
        return (None, "")


def _emu_to_inches(emu: int | None) -> float:
    if emu is None:
        return 0.0
    return round(emu / 914400.0, 2)


def _shape_type_repr(shape_type) -> str:
    return str(shape_type) if shape_type is not None else "None"


def collect_shape_records(slide, slide_number: int) -> list[ShapeRecord]:
    """Walk every shape on ``slide`` (expanding GROUPs), returning records."""
    records: list[ShapeRecord] = []
    counter = 0

    def walk(shape_iter: Iterable, prefix: str = "") -> None:
        nonlocal counter
        for shape in shape_iter:
            if shape.shape_type == 6:  # GROUP
                walk(shape.shapes, prefix=f"{shape.name}/")
                continue
            text = ""
            runs: list[RunRecord] = []
            if shape.has_text_frame:
                text = shape.text_frame.text or ""
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        target, _ = extract_run_hyperlink(run)
                        runs.append(RunRecord(text=run.text or "", hyperlink=target))
            shape_target, _ = extract_shape_hyperlink(shape)
            records.append(
                ShapeRecord(
                    slide_number=slide_number,
                    index=counter,
                    name=f"{prefix}{shape.name}",
                    shape_type=_shape_type_repr(shape.shape_type),
                    left_in=_emu_to_inches(shape.left),
                    top_in=_emu_to_inches(shape.top),
                    width_in=_emu_to_inches(shape.width),
                    height_in=_emu_to_inches(shape.height),
                    text=text,
                    shape_hyperlink=shape_target,
                    runs=runs,
                )
            )
            counter += 1

    walk(slide.shapes)
    return records


def _ext_of(target: str) -> str:
    suffix = Path(target).suffix.lower()
    return suffix if suffix else "(no-ext)"


def _collapse_ws(text: str) -> str:
    """Collapse any run of whitespace (incl. newlines/tabs) to a single space."""
    return " ".join(text.split())


def render_summary(
    filename: str,
    total_slides: int,
    all_records: list[ShapeRecord],
    grams_by_slide: dict[int, list[GramPlaceholder]] | None = None,
    framing_slides: list[int] | None = None,
) -> str:
    """Render section 1 of the report.

    ``grams_by_slide`` and ``framing_slides`` are optional so existing
    callers (and unit tests) that pre-date the gram-aware refactor still
    work; when omitted the gram-centric counts are simply skipped.
    """
    targets_by_ext_live: dict[str, int] = defaultdict(int)
    shape_level_live = 0
    shape_level_vestigial = 0
    text_run = 0
    by_slide: dict[int, int] = defaultdict(int)
    for rec in all_records:
        by_slide[rec.slide_number] += 1
        if rec.shape_hyperlink:
            if rec.shape_hyperlink.lower().startswith("file:///"):
                shape_level_vestigial += 1
            else:
                shape_level_live += 1
                targets_by_ext_live[_ext_of(rec.shape_hyperlink)] += 1
        for run in rec.runs:
            if run.hyperlink:
                text_run += 1
                targets_by_ext_live[_ext_of(run.hyperlink)] += 1

    deviating: list[int] = []
    # Slide 1 (welcome) and the last slide (exit) are framing slides — they
    # carry only a title + subtitle text box, not the gram grid, so the
    # deviation check excludes them.
    for slide_num in range(2, total_slides):
        count = by_slide.get(slide_num, 0)
        if abs(count - EXPECTED_SHAPES_PER_CONTENT_SLIDE) > SHAPE_DEVIATION_TOLERANCE:
            deviating.append(slide_num)

    lines = ["=== Section 1: Summary ===", f"Filename: {filename}", f"Total slides: {total_slides}"]
    if framing_slides is not None:
        framing_label = ", ".join(str(n) for n in framing_slides) or "none"
        lines.append(f"Framing slides (skipped): {framing_label}")
    if grams_by_slide is not None:
        all_grams = [g for grams in grams_by_slide.values() for g in grams]
        no_analysis = sum(1 for g in all_grams if not g.png_href)
        no_glc = sum(1 for g in all_grams if not g.glc_links)
        lines.append(f"Total grams extracted: {len(all_grams)}")
        lines.append(f"  grams missing analysis sheet href: {no_analysis}")
        lines.append(f"  grams missing .glc links: {no_glc}")
    lines.append("Hyperlink target extensions (live only):")
    for ext in sorted(targets_by_ext_live):
        lines.append(f"  {ext}: {targets_by_ext_live[ext]}")
    lines.append(f"Shape-level hyperlinks (live): {shape_level_live}")
    lines.append(f"Shape-level hyperlinks (vestigial absolute file:///): {shape_level_vestigial}")
    lines.append(f"Text-run hyperlinks: {text_run}")
    if deviating:
        lines.append(f"Slides flagged (deviating shape count): {', '.join(str(n) for n in deviating)}")
    else:
        lines.append("Slides flagged (deviating shape count): none")
    return "\n".join(lines) + "\n"


def render_per_gram(
    grams_by_slide: dict[int, list[GramPlaceholder]],
    slides_filter: list[int] | None,
) -> str:
    """Render section 2 (default view): grams as the extractor will see them."""
    lines = ["=== Section 2: Per-gram ==="]
    for slide_num in sorted(grams_by_slide):
        if slides_filter is not None and slide_num not in slides_filter:
            continue
        grams = grams_by_slide[slide_num]
        if not grams:
            lines.append(f"-- Slide {slide_num} (no grams — framing or anomalous) --")
            continue
        # Hide grams that resolved no .glc links — typically the
        # "student" half of mixed student/instructor decks, where the
        # gram header exists but its lofar links don't resolve. Counted
        # in the slide header so the omission stays visible.
        with_links = [g for g in grams if g.glc_links]
        hidden = len(grams) - len(with_links)
        if not with_links:
            lines.append(
                f"-- Slide {slide_num} (no grams with lofars; "
                f"{hidden} header-only gram(s) hidden) --"
            )
            continue
        suffix = f"; {hidden} header-only hidden" if hidden else ""
        lines.append(f"-- Slide {slide_num} ({len(with_links)} grams{suffix}) --")
        # Grams already arrive sorted by gram_id from extract_grams_from_slide;
        # CSV and report agree on order.
        for idx, g in enumerate(with_links, start=1):
            title_bits = [g.gram_id or "(no id)"]
            if g.vessel_name:
                title_bits.append(_collapse_ws(g.vessel_name))
            lines.append(f"  [{idx}] {': '.join(title_bits)}")
            lines.append(f"      analysis: {g.png_href or '(MISSING)'}")
            if g.glc_links:
                for j, link in enumerate(g.glc_links, start=1):
                    text = _collapse_ws(link.display_text)
                    lines.append(f"      glc[{j}]: {link.href}  text={text!r}")
            else:
                lines.append("      glc: (MISSING)")
    return "\n".join(lines) + "\n"


def render_per_slide(all_records: list[ShapeRecord], slides_filter: list[int] | None) -> str:
    """Render section 2 of the report."""
    lines = ["=== Section 3: Per-slide (verbose) ==="]
    by_slide: dict[int, list[ShapeRecord]] = defaultdict(list)
    for rec in all_records:
        if slides_filter is not None and rec.slide_number not in slides_filter:
            continue
        by_slide[rec.slide_number].append(rec)
    for slide_num in sorted(by_slide):
        recs = by_slide[slide_num]
        lines.append(f"-- Slide {slide_num} (shapes: {len(recs)}) --")
        for rec in recs:
            text = rec.text.replace("\n", " ")
            if len(text) > TEXT_TRUNCATION:
                text = text[:TEXT_TRUNCATION] + "..."
            lines.append(
                f"  [{rec.index}] name={rec.name} type={rec.shape_type} "
                f"pos=({rec.left_in:.2f},{rec.top_in:.2f}) "
                f"size=({rec.width_in:.2f}x{rec.height_in:.2f}) "
                f"text={text!r} shape_hyperlink={rec.shape_hyperlink}"
            )
            for j, run in enumerate(rec.runs):
                lines.append(f"      run[{j}]: text={run.text!r} hyperlink={run.hyperlink}")
    return "\n".join(lines) + "\n"


def render_hyperlinks(all_records: list[ShapeRecord]) -> str:
    """Render section 3 of the report."""
    by_ext: dict[str, list[tuple[str, str, int, str]]] = defaultdict(list)
    seen: set[tuple[str, str, int, str]] = set()
    for rec in all_records:
        if rec.shape_hyperlink:
            entry = (rec.shape_hyperlink, "shape-level", rec.slide_number, rec.name)
            if entry not in seen:
                seen.add(entry)
                by_ext[_ext_of(rec.shape_hyperlink)].append(entry)
        for run in rec.runs:
            if run.hyperlink:
                entry = (run.hyperlink, "text-run", rec.slide_number, rec.name)
                if entry not in seen:
                    seen.add(entry)
                    by_ext[_ext_of(run.hyperlink)].append(entry)

    lines = ["=== Section 4: Hyperlink targets (verbose, raw) ==="]
    for ext in sorted(by_ext):
        lines.append(f"-- {ext} --")
        for target, kind, slide_num, shape_name in sorted(by_ext[ext]):
            lines.append(f"  {target} [{kind}] slide={slide_num} shape={shape_name}")
    return "\n".join(lines) + "\n"


def _parse_slides_filter(value: str | None) -> list[int] | None:
    if not value:
        return None
    return [int(v) for v in value.split(",") if v.strip()]


def main(argv: Iterable[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description="Produce a structural report for a PPTX")
    parser.add_argument("--input", required=True, type=Path)
    parser.add_argument("--out", required=False, type=Path, default=None)
    parser.add_argument("--slides", required=False, default=None,
                        help="Comma-separated slide numbers to include in per-gram and verbose sections")
    parser.add_argument("--verbose", action="store_true",
                        help="Append per-shape and raw-hyperlink forensic sections")
    args = parser.parse_args(list(argv) if argv is not None else None)

    setup_logging(Path("introspect.log"))
    LOGGER.info("Opening %s", args.input)
    try:
        prs = Presentation(args.input)
    except Exception as exc:
        LOGGER.error("Cannot open PPTX %s: %s", args.input, exc)
        return 1

    all_records: list[ShapeRecord] = []
    grams_by_slide: dict[int, list[GramPlaceholder]] = {}
    framing_slides: list[int] = []
    total_slides = 0
    for i, slide in enumerate(prs.slides, start=1):
        total_slides = i
        all_records.extend(collect_shape_records(slide, slide_number=i))
        # Synthesize records for SmartArt-embedded hyperlinks so the
        # raw dump (Section 4) shows them alongside slide-level links.
        for j, (_text, href) in enumerate(_slide_diagram_hyperlinks(slide)):
            all_records.append(ShapeRecord(
                slide_number=i,
                index=10_000 + j,
                name="<smartart-diagram>",
                shape_type="diagram-node",
                left_in=0.0, top_in=0.0, width_in=0.0, height_in=0.0,
                text="",
                shape_hyperlink=href,
                runs=[],
            ))
        if is_framing_slide(slide):
            framing_slides.append(i)
            grams_by_slide[i] = []
        else:
            grams_by_slide[i] = extract_grams_from_slide(
                slide, slide_num=i,
                content_root=args.input.parent,
                source_dir=args.input.parent,
            )

    slides_filter = _parse_slides_filter(args.slides)
    sections = [
        render_summary(args.input.name, total_slides, all_records, grams_by_slide, framing_slides),
        render_per_gram(grams_by_slide, slides_filter),
    ]
    if args.verbose:
        sections.append(render_per_slide(all_records, slides_filter))
        sections.append(render_hyperlinks(all_records))
    report = "\n".join(sections)

    if args.out is not None:
        args.out.parent.mkdir(parents=True, exist_ok=True)
        args.out.write_text(report, encoding="utf-8")
        LOGGER.info("Wrote report to %s", args.out)
    else:
        sys.stdout.write(report)
    return 0


if __name__ == "__main__":
    rc = main()
    # Preserve CLI exit codes when invoked as a script, but stay silent
    # when invoked from an interactive REPL via runpy.run_path —
    # ``sys.exit`` would otherwise kill the interpreter and break the
    # up-arrow iteration loop. ``sys.ps1`` is only defined in
    # interactive sessions.
    if rc and not hasattr(sys, "ps1"):
        sys.exit(rc)
